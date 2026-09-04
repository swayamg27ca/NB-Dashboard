# ==============================================================================
# STREAMLIT LIFE INSURANCE EXECUTIVE DASHBOARD - BOARDROOM EDITION
# ==============================================================================
import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import re
import os
import io
import base64
from pathlib import Path
from datetime import datetime

# REQUIRED PLOTLY IMPORTS
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots

# PPTX IMPORT WITH FALLBACK HANDLING
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# REPORTLAB (PDF) IMPORT WITH FALLBACK HANDLING
try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, KeepTogether
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    HAS_REPORTLAB = True
except ImportError:
    HAS_REPORTLAB = False

# Define base directory
BASE_DIR = Path(__file__).parent if "__file__" in dir() else Path.cwd()

# Page Configuration
st.set_page_config(
    page_title="CreditAccess Life New Business Dashboard",
    page_icon=str(BASE_DIR / "favicon.png") if (BASE_DIR / "favicon.png").exists() else "CALI",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ==============================================================================
# 🔒 SECURE LOGIN PASSWORD SYSTEM
# ==============================================================================
ENABLE_LOGIN = True
DASHBOARD_PASSWORD = "Sairam@2026"

def get_target_password():
    try:
        if hasattr(st, "secrets") and "DASHBOARD_PASSWORD" in st.secrets:
            return st.secrets["DASHBOARD_PASSWORD"]
    except Exception:
        pass
    return DASHBOARD_PASSWORD

def check_password():
    if not ENABLE_LOGIN:
        return True

    if "authenticated" not in st.session_state:
        st.session_state["authenticated"] = False

    if not st.session_state["authenticated"]:
        target_pwd = get_target_password()
        st.markdown("<br><br>", unsafe_allow_html=True)
        col_sec1, col_sec2, col_sec3 = st.columns([1, 1.5, 1])
        with col_sec2:
            st.markdown("""
                <div style='background-color: #1e3a5f; padding: 28px; border-radius: 12px; border: 2px solid #e63946; text-align: center; box-shadow: 0px 6px 20px rgba(0,0,0,0.5);'>
                    <span style='background: #e63946; color: white; padding: 4px 12px; border-radius: 4px; font-size: 12px; font-weight: bold; letter-spacing: 1px;'>SECURE ACCESS</span>
                    <h2 style='color:#ffffff; margin-top: 12px; margin-bottom: 6px; font-weight: 800;'>Executive Dashboard Login</h2>
                    <p style='color:#f1faee; font-size: 16px; margin: 0;'>Suraksha aur Samruddhi | CreditAccess Life Analytics</p>
                </div>
            """, unsafe_allow_html=True)
            st.markdown("<br>", unsafe_allow_html=True)
            
            pwd_input = st.text_input("Enter Dashboard Password", type="password", key="login_pwd")
            if st.button("Login to Dashboard", use_container_width=True):
                if pwd_input == target_pwd:
                    st.session_state["authenticated"] = True
                    st.success("Authentication Successful")
                    st.rerun()
                else:
                    st.error("Access Denied: Incorrect Password")
        return False
    return True

if not check_password():
    st.stop()

# Helper for Base64 image rendering
def get_image_base64(path):
    try:
        with open(path, "rb") as img_file:
            return base64.b64encode(img_file.read()).decode()
    except Exception:
        return None

# Helper for converting DataFrames to Excel bytes
def convert_df_to_excel(df, sheet_name="Sheet1"):
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        df.to_excel(writer, sheet_name=sheet_name, index=False)
    output.seek(0)
    return output.getvalue()

# ==============================================================================
# 🎨 BOARDROOM REPORT GENERATORS (.PPTX & REAL PDF)
# ==============================================================================
def generate_boardroom_pptx(data, month_name):
    """Generates a complete multi-slide executive PowerPoint presentation deck (.pptx) covering all tabs"""
    if not HAS_PPTX:
        return None

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_header(slide, title_text):
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(30, 58, 95)

    def add_table_slide(title, df, max_r=10):
        s = prs.slides.add_slide(blank_layout)
        add_header(s, title)
        
        sub_df = df.head(max_r).copy()
        rows, cols = len(sub_df) + 1, len(sub_df.columns)
        table_shape = s.shapes.add_table(rows, cols, Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.2))
        table = table_shape.table

        for c_idx, col_name in enumerate(sub_df.columns):
            cell = table.cell(0, c_idx)
            cell.text = str(col_name)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(230, 57, 70)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.font.bold = True
                paragraph.font.size = Pt(11)

        for r_idx, (_, row) in enumerate(sub_df.iterrows()):
            r = r_idx + 1
            for c_idx, val in enumerate(row):
                cell = table.cell(r, c_idx)
                if isinstance(val, (int, float, np.number)):
                    cell.text = f"{val:,.2f}" if isinstance(val, float) else f"{val:,.0f}"
                else:
                    cell.text = str(val)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(10)
                    paragraph.font.color.rgb = RGBColor(30, 58, 95)

    # --- SLIDE 1: COVER ---
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = RGBColor(30, 58, 95)
    bg1.line.color.rgb = RGBColor(30, 58, 95)

    accent_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(0.15), Inches(4.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = RGBColor(230, 57, 70)
    accent_bar.line.color.rgb = RGBColor(230, 57, 70)

    txBox = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "CreditAccess Life New Business Dashboard"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    p2 = tf.add_paragraph()
    p2.text = f"Complete Boardroom Performance Deck | {month_name.replace('_', ' ')}"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(241, 250, 238)

    # --- SLIDE 2: KPI & EXECUTIVE BRIEFING ---
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, f"Executive Market Overview – {month_name.replace('_', ' ')}")

    tot_fyp = data["df_summary"]['FYP (in Cr.)'].sum() if not data["df_summary"].empty else 0.0
    tot_lives = data["df_lives_summary"]['Lives Covered'].sum() if not data["df_lives_summary"].empty else 0.0
    cali_row = data["final_growth_df"][data["final_growth_df"]['INSURER'].str.upper().isin(['CALI', 'CREDITACCESS', 'CREDITACCESS LIFE'])]
    cali_rank = f"#{int(cali_row['Rank'].values[0])}" if not cali_row.empty else "N/A"
    cali_fyp = f"Rs. {cali_row[data['col_current_ytd']].values[0]:,.2f} Cr" if not cali_row.empty else "Rs. 0.00"

    metrics = [
        ("Industry Total FYP", f"Rs. {tot_fyp:,.2f} Cr"),
        ("Total Lives Covered", f"{tot_lives:,.0f}"),
        ("CALI Market Rank", cali_rank),
        ("CALI YTD FYP", cali_fyp)
    ]

    for idx, (label, val) in enumerate(metrics):
        x_pos = Inches(0.8 + idx * 3.0)
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(1.5), Inches(2.7), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(248, 249, 250)
        card.line.color.rgb = RGBColor(230, 57, 70)
        
        tf_card = card.text_frame
        tf_card.word_wrap = True
        p_c1 = tf_card.paragraphs[0]
        p_c1.text = label.upper()
        p_c1.font.size = Pt(11)
        p_c1.font.color.rgb = RGBColor(100, 100, 100)
        p_c1.font.bold = True
        
        p_c2 = tf_card.add_paragraph()
        p_c2.text = val
        p_c2.font.size = Pt(18)
        p_c2.font.bold = True
        p_c2.font.color.rgb = RGBColor(30, 58, 95)

    add_table_slide("Tab 1: Segment FYP Summary", data["df_summary"])
    add_table_slide("Tab 1: Complete FYP Matrix (YTD Cr)", data["fyp_matrix"], max_r=12)
    add_table_slide("Tab 2: Top Insurers YoY Growth Performance", data["final_growth_df"], max_r=12)
    add_table_slide("Tab 2: Portfolio Split (% Proportion)", data["df_ms"], max_r=12)
    add_table_slide("Tab 5: Premium Growth Matrix (%)", data["growth_combined"], max_r=12)
    add_table_slide("Tab 6: Segment Lives Covered Summary", data["df_lives_summary"])
    add_table_slide("Tab 6: Master Policy Schemes (MPH) Matrix", data["mph_matrix"], max_r=12)
    add_table_slide("Tab 7: Average Premium Per Life (Rs)", data["df_avg_life"], max_r=12)
    add_table_slide("Tab 8: Executive Annexure", data["df_ann"], max_r=12)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


def generate_boardroom_pdf(data, month_name):
    """Generates a complete multi-page executive PDF report (.pdf) covering all 10 tabs"""
    if not HAS_REPORTLAB:
        return None

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, rightMargin=36, leftMargin=36, topMargin=36, bottomMargin=36)
    story = []

    styles = getSampleStyleSheet()
    
    title_style = ParagraphStyle(
        'DocTitle', parent=styles['Normal'],
        fontName='Helvetica-Bold', fontSize=20, textColor=colors.HexColor("#1e3a5f"),
        spaceAfter=4
    )
    sub_style = ParagraphStyle(
        'DocSubTitle', parent=styles['Normal'],
        fontName='Helvetica', fontSize=12, textColor=colors.HexColor("#e63946"),
        spaceAfter=12
    )
    h2_style = ParagraphStyle(
        'SectionH2', parent=styles['Normal'],
        fontName='Helvetica-Bold', fontSize=13, textColor=colors.HexColor("#1e3a5f"),
        spaceBefore=12, spaceAfter=6
    )
    body_style = ParagraphStyle(
        'BodyDark', parent=styles['Normal'],
        fontName='Helvetica', fontSize=9, textColor=colors.HexColor("#333333"),
        leading=13
    )
    comment_style = ParagraphStyle(
        'AnalysisComment', parent=styles['Normal'],
        fontName='Helvetica-Oblique', fontSize=9, textColor=colors.HexColor("#1e3a5f"),
        backColor=colors.HexColor("#f0f4f8"), borderColor=colors.HexColor("#e63946"),
        borderWidth=1, borderPadding=6, spaceBefore=6, spaceAfter=10, leading=13
    )

    def make_pdf_table(df, col_widths=None, header_bg="#e63946", max_r=25):
        sub_df = df.head(max_r).copy()
        table_data = [[Paragraph(f"<b>{col}</b>", ParagraphStyle('TH', parent=body_style, textColor=colors.white, alignment=1, fontSize=8)) for col in sub_df.columns]]
        
        for _, row in sub_df.iterrows():
            r_cells = []
            for val in row:
                if isinstance(val, float):
                    txt = f"{val:,.2f}"
                elif isinstance(val, (int, np.integer)):
                    txt = f"{val:,.0f}"
                else:
                    txt = str(val)
                r_cells.append(Paragraph(txt, ParagraphStyle('TD', parent=body_style, alignment=1, fontSize=8)))
            table_data.append(r_cells)

        t = Table(table_data, colWidths=col_widths)
        t.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.HexColor(header_bg)),
            ('ALIGN', (0,0), (-1,-1), 'CENTER'),
            ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
            ('BOTTOMPADDING', (0,0), (-1,0), 4),
            ('TOPPADDING', (0,0), (-1,0), 4),
            ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#d0d0d0')),
            ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.white, colors.HexColor('#f8f9fa')])
        ]))
        return t

    # --- COVER HEADER & KPI CARDS ---
    story.append(Paragraph("CreditAccess Life New Business Dashboard", title_style))
    story.append(Paragraph(f"Complete Boardroom Analytics Report | {month_name.replace('_', ' ')}", sub_style))
    
    tot_fyp = data["df_summary"]['FYP (in Cr.)'].sum() if not data["df_summary"].empty else 0.0
    tot_lives = data["df_lives_summary"]['Lives Covered'].sum() if not data["df_lives_summary"].empty else 0.0
    cali_row = data["final_growth_df"][data["final_growth_df"]['INSURER'].str.upper().isin(['CALI', 'CREDITACCESS', 'CREDITACCESS LIFE'])]
    cali_rank = f"#{int(cali_row['Rank'].values[0])}" if not cali_row.empty else "N/A"
    cali_fyp = f"Rs. {cali_row[data['col_current_ytd']].values[0]:,.2f} Cr" if not cali_row.empty else "Rs. 0.00"
    top_co = data["final_growth_df"].iloc[0]['INSURER'] if not data["final_growth_df"].empty else "N/A"

    kpi_data = [
        [Paragraph("<b>INDUSTRY TOTAL FYP</b>", body_style), Paragraph("<b>TOTAL LIVES COVERED</b>", body_style), Paragraph("<b>CALI MARKET RANK</b>", body_style), Paragraph("<b>CALI YTD FYP</b>", body_style)],
        [Paragraph(f"<b>Rs. {tot_fyp:,.2f} Cr</b>", sub_style), Paragraph(f"<b>{tot_lives:,.0f}</b>", sub_style), Paragraph(f"<b>{cali_rank}</b>", sub_style), Paragraph(f"<b>{cali_fyp}</b>", sub_style)]
    ]
    t_kpi = Table(kpi_data, colWidths=[130, 130, 130, 130])
    t_kpi.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#f8f9fa')),
        ('BOX', (0,0), (-1,-1), 1, colors.HexColor('#e63946')),
        ('INNERGRID', (0,0), (-1,-1), 0.5, colors.HexColor('#e0e0e0')),
        ('ALIGN', (0,0), (-1,-1), 'CENTER')
    ]))
    story.append(t_kpi)
    story.append(Spacer(1, 10))

    story.append(Paragraph(f"<b>Executive Summary & Strategic Analysis:</b> Industry First Year Premium stands at <b>Rs. {tot_fyp:,.2f} Cr</b> with <b>{top_co}</b> leading private sector volume. CreditAccess Life (CALI) is positioned at <b>Rank {cali_rank}</b> with <b>{cali_fyp}</b>. Group Single Premium (GSP) and Individual Non-Single (INSP) continue to represent key segment growth drivers.", comment_style))

    # --- PAGE 1: TAB 1 MARKET OVERVIEW & FYP MATRIX ---
    story.append(Paragraph("1. Industry FYP Segment Summary", h2_style))
    story.append(make_pdf_table(data["df_summary"], col_widths=[180, 170, 170], header_bg="#1e3a5f"))
    story.append(Spacer(1, 10))

    story.append(Paragraph("2. Complete FYP Matrix (YTD Cr)", h2_style))
    story.append(make_pdf_table(data["fyp_matrix"], col_widths=[110, 68, 68, 68, 68, 68, 68], header_bg="#e63946", max_r=15))
    story.append(PageBreak())

    # --- PAGE 2: TAB 2 YoY GROWTH & PORTFOLIO SPLIT ---
    story.append(Paragraph("3. Private Insurers YoY Growth Leaderboard", h2_style))
    story.append(make_pdf_table(data["final_growth_df"], col_widths=[40, 130, 115, 115, 120], header_bg="#e63946", max_r=15))
    story.append(Paragraph("<b>Growth Comment:</b> Performance reflects strong demand in group channels. CALI's growth trajectories are tracked relative to peer benchmarks.", body_style))
    story.append(Spacer(1, 10))

    story.append(Paragraph("4. Individual vs Group Portfolio Split", h2_style))
    story.append(make_pdf_table(data["df_ms"], col_widths=[40, 110, 90, 90, 95, 95], header_bg="#1e3a5f", max_r=15))
    story.append(PageBreak())

    # --- PAGE 3: TAB 5 GROWTH MATRIX & RATES ---
    story.append(Paragraph("5. Premium Growth Matrix (%)", h2_style))
    story.append(make_pdf_table(data["growth_combined"], col_widths=[110, 68, 68, 68, 68, 68, 68], header_bg="#1e3a5f", max_r=15))
    story.append(Spacer(1, 10))

    story.append(Paragraph("6. Premium Rates (Premium / SA × 1000)", h2_style))
    story.append(make_pdf_table(data["df_rate_all"], col_widths=[110, 68, 68, 68, 68, 68, 68], header_bg="#e63946", max_r=15))
    story.append(PageBreak())

    # --- PAGE 4: TAB 6 LIVES & MPH MATRICES ---
    story.append(Paragraph("7. Segment-Wise Summary of Lives Covered", h2_style))
    story.append(make_pdf_table(data["df_lives_summary"], col_widths=[180, 170, 170], header_bg="#1e3a5f"))
    story.append(Spacer(1, 10))

    story.append(Paragraph("8. Master Policy Schemes (MPH) Matrix", h2_style))
    story.append(make_pdf_table(data["mph_matrix"], col_widths=[110, 68, 68, 68, 68, 68, 68], header_bg="#e63946", max_r=15))
    story.append(Spacer(1, 10))

    story.append(Paragraph("9. Complete Lives Covered Matrix", h2_style))
    story.append(make_pdf_table(data["df_lives_clean"], col_widths=[110, 68, 68, 68, 68, 68, 68], header_bg="#1e3a5f", max_r=15))
    story.append(PageBreak())

    # --- PAGE 5: TAB 7 & 8 AVERAGE METRICS & ANNEXURE ---
    story.append(Paragraph("10. Average Premium Per Life (Rs)", h2_style))
    story.append(make_pdf_table(data["df_avg_life"], col_widths=[110, 68, 68, 68, 68, 68, 68], header_bg="#e63946", max_r=12))
    story.append(Spacer(1, 10))

    story.append(Paragraph("11. Executive Annexure", h2_style))
    story.append(make_pdf_table(data["df_ann"], col_widths=[160, 120, 120, 120], header_bg="#1e3a5f", max_r=15))

    doc.build(story)
    buffer.seek(0)
    return buffer.getvalue()


# ==============================================================================
# CORPORATE CSS STYLING & MULTISELECT BADGES VISIBILITY FIX
# ==============================================================================
st.markdown("""
<style>
    /* 1. HIDE TOP HEADER BAR & UTILIZE TOP SPACE EFFICIENTLY */
    [data-testid="stHeader"] {
        background-color: transparent !important;
        height: 0px !important;
    }

    /* 2. OPTIMIZED CONTAINER PADDING */
    .main .block-container {
        padding-top: 0.8rem !important;
        padding-bottom: 1rem !important;
        padding-left: 1.5rem !important;
        padding-right: 1.5rem !important;
        max-width: 100% !important;
    }

    /* 3. LIGHTER GRADIENT BACKGROUND */
    .stApp {
        background: linear-gradient(135deg, #2b4c7e 0%, #4a6fa5 30%, #c0392b 80%, #e63946 100%) !important;
        color: #ffffff !important;
    }
    
    /* COMPACT SIDEBAR STYLING */
    [data-testid="stSidebar"] {
        background-color: #1e3a5f !important;
        border-right: 2px solid #e63946 !important;
        padding-top: 1rem !important;
    }
    [data-testid="stSidebar"] h1, [data-testid="stSidebar"] h2, [data-testid="stSidebar"] h3 {
        font-size: 17px !important;
        color: #ffffff !important;
        font-weight: 700 !important;
    }
    [data-testid="stSidebar"] p, [data-testid="stSidebar"] span, [data-testid="stSidebar"] label {
        font-size: 15px !important;
        color: #f1faee !important;
    }

    /* Global Font Size Enforcements */
    html, body, [class*="css"]  {
        font-size: 22px !important;
        color: #ffffff !important;
    }
    
    h1 {
        font-size: 42px !important;
        color: #ffffff !important;
        text-shadow: 2px 2px 4px rgba(0,0,0,0.4);
        font-weight: 800 !important;
        margin: 0 !important;
    }
    h2 { font-size: 28px !important; color: #ffffff !important; }
    h3, h4, h5 { font-size: 24px !important; color: #ffffff !important; }
    p, span, label { font-size: 20px !important; color: #ffffff !important; }
    
    /* FIX SELECTBOXES, MULTISELECTS & DROPDOWN MENUS */
    div[data-baseweb="select"], div[data-baseweb="select"] *,
    div[data-baseweb="input"], div[data-baseweb="input"] *,
    div[data-baseweb="popover"], div[data-baseweb="popover"] *,
    div[data-baseweb="menu"], div[data-baseweb="menu"] *,
    ul[role="listbox"], ul[role="listbox"] *,
    li[role="option"], li[role="option"] * {
        background-color: #ffffff !important;
        color: #1e3a5f !important;
        font-weight: 600 !important;
    }

    /* MULTISELECT TAG BADGES FIX (DARK BLUE PILL WITH BOLD WHITE TEXT) */
    div[data-baseweb="select"] span[data-baseweb="tag"] {
        background-color: #1e3a5f !important;
        border: 1px solid #e63946 !important;
        border-radius: 6px !important;
        padding: 3px 10px !important;
        margin: 2px !important;
    }
    div[data-baseweb="select"] span[data-baseweb="tag"] *,
    div[data-baseweb="select"] span[data-baseweb="tag"] span,
    div[data-baseweb="select"] span[data-baseweb="tag"] div {
        background-color: transparent !important;
        color: #ffffff !important;
        font-size: 14px !important;
        font-weight: 800 !important;
        opacity: 1 !important;
    }

    /* FIX FILE UPLOADER DROPZONE VISIBILITY */
    [data-testid="stFileUploader"] {
        background-color: #ffffff !important;
        border-radius: 10px !important;
        padding: 12px !important;
        border: 2px dashed #e63946 !important;
    }
    [data-testid="stFileUploader"] *,
    [data-testid="stFileUploader"] label,
    [data-testid="stFileUploader"] span,
    [data-testid="stFileUploader"] p,
    [data-testid="stFileUploader"] small,
    [data-testid="stFileUploader"] div {
        color: #1e3a5f !important;
        font-weight: 700 !important;
        opacity: 1 !important;
    }
    [data-testid="stFileUploader"] button {
        background-color: #e63946 !important;
        border: none !important;
    }
    [data-testid="stFileUploader"] button * {
        color: #ffffff !important;
        font-weight: bold !important;
    }

    /* RESET PLOTLY MODEBAR TOOLBAR BUTTONS */
    .js-plotly-plot .plotly .modebar-container, .modebar-btn {
        background: transparent !important;
        box-shadow: none !important;
        border: none !important;
    }

    /* TARGET STREAMLIT BUTTONS */
    div.stButton > button, div[data-testid="stDownloadButton"] > button {
        background-color: #e63946 !important;
        color: #ffffff !important;
        border: none !important;
        border-radius: 8px !important;
        padding: 8px 16px !important;
        box-shadow: 0px 3px 8px rgba(0,0,0,0.25) !important;
        transition: all 0.3s ease !important;
    }
    div.stButton > button *, div[data-testid="stDownloadButton"] > button * {
        color: #ffffff !important;
        font-weight: bold !important;
        font-size: 16px !important;
    }
    div.stButton > button:hover, div[data-testid="stDownloadButton"] > button:hover {
        background-color: #c1121f !important;
        transform: translateY(-2px) !important;
    }

    /* Metric Card Styling & Hover Animation */
    [data-testid="stMetric"] {
        background: rgba(30, 58, 95, 0.85) !important;
        backdrop-filter: blur(8px) !important;
        border: 1px solid rgba(255, 255, 255, 0.15) !important;
        border-left: 5px solid #e63946 !important;
        padding: 12px !important;
        border-radius: 10px !important;
        transition: transform 0.25s ease, box-shadow 0.25s ease !important;
        box-shadow: 0px 4px 12px rgba(0,0,0,0.25) !important;
    }
    [data-testid="stMetric"]:hover {
        transform: translateY(-4px) !important;
        box-shadow: 0px 8px 20px rgba(230, 57, 70, 0.4) !important;
    }
    [data-testid="stMetricValue"] { font-size: 28px !important; color: #ffffff !important; font-weight: 800 !important; }
    [data-testid="stMetricLabel"] { font-size: 16px !important; color: #f1faee !important; text-transform: uppercase; letter-spacing: 0.5px; }

    /* DATAFRAME STYLING & STRICT CENTER ALIGNMENT */
    div[data-testid="stDataFrame"] {
        background-color: rgba(30, 58, 95, 0.85) !important;
        border-radius: 12px;
        padding: 8px;
    }
    
    div[data-testid="stDataFrame"] th,
    div[data-testid="stDataFrame"] td,
    div[data-testid="stDataFrame"] [role="columnheader"],
    div[data-testid="stDataFrame"] [role="gridcell"],
    div[data-testid="stDataFrame"] [data-testid="stTable"] * {
        text-align: center !important;
        justify-content: center !important;
    }

    div[data-testid="stDataFrame"] th {
        background-color: #e63946 !important;
        color: #ffffff !important;
        font-size: 19px !important;
        text-align: center !important;
    }
    div[data-testid="stDataFrame"] td {
        color: #ffffff !important;
        font-size: 18px !important;
        text-align: center !important;
    }

    /* HIDE UGLY FLOATING WHITE DATAFRAME TOOLBAR */
    [data-testid="stDataFrameToolbar"],
    [data-testid="stDataFrameToolbar"] *,
    div[data-testid="stDataFrame"] > div:first-child > div:nth-child(2) {
        display: none !important;
        visibility: hidden !important;
        opacity: 0 !important;
        height: 0px !important;
        width: 0px !important;
        pointer-events: none !important;
    }
    
    /* FORCE SOLID WHITE TEXT FOR ALL TABS */
    .stTabs { margin-top: 10px !important; }
    .stTabs [data-baseweb="tab"] {
        background-color: rgba(15, 30, 50, 0.75) !important;
        border-radius: 6px 6px 0 0;
        padding: 10px 16px;
    }
    .stTabs [data-baseweb="tab"] * { color: #ffffff !important; opacity: 1 !important; font-weight: 600 !important; font-size: 18px !important; }
    .stTabs [aria-selected="true"] { background-color: #e63946 !important; }
    .stTabs [aria-selected="true"] * { color: #ffffff !important; font-weight: 800 !important; }
</style>
""", unsafe_allow_html=True)

# ==============================================================================
# MAPPINGS & CONSTANTS
# ==============================================================================
segment_mapping = {
    'Individual Single Premium': 'ISP',
    'Individual Non-Single Premium': 'INSP',
    'Group Single Premium': 'GSP',
    'Group Non-Single Premium': 'GNSP',
    'Group Yearly Renewable Premium': 'GYRP'
}

company_keywords = [
    ('sbi', 'SBI Life'), ('hdfc', 'HDFC Life'), ('icici', 'ICICI Prudential'),
    ('bajaj', 'Bajaj Allianz'), ('axis', 'Axis Max Life'), ('tata', 'Tata AIA'),
    ('aditya', 'Aditya Birla'), ('kotak', 'Kotak Mahindra'), ('star union', 'Star Union Dai-ichi'),
    ('canara', 'Canara HSBC'), ('pnb', 'PNB MetLife'), ('indiafirst', 'IndiaFirst Life'),
    ('shriram', 'Shriram'), ('pramerica', 'Pramerica Life'), ('digit', 'Go Digit'),
    ('generali', 'Generali Central'), ('indusind', 'IndusInd Nippon'), ('ageas', 'Ageas Federal'),
    ('bharti', 'Bharti AXA'), ('aviva', 'Aviva'), ('bandhan', 'Bandhan'),
    ('creditaccess', 'CALI'), ('edelweiss', 'Edelweiss'), ('acko', 'Acko')
]

MONTH_CALENDAR_ORDER = {
    'january': 1, 'jan': 1, 'february': 2, 'feb': 2, 'march': 3, 'mar': 3,
    'april': 4, 'apr': 4, 'may': 5, 'june': 6, 'jun': 6, 'july': 7, 'jul': 7,
    'august': 8, 'aug': 8, 'september': 9, 'sep': 9, 'sept': 9, 'october': 10, 'oct': 10,
    'november': 11, 'nov': 11, 'december': 12, 'dec': 12
}

def parse_month_year(filename):
    s = str(filename).lower()
    month_val, year_val = 99, 9999
    for m_name, m_num in MONTH_CALENDAR_ORDER.items():
        if m_name in s:
            month_val = m_num
            break
    y_match = re.search(r'20\d{2}', s)
    if y_match:
        year_val = int(y_match.group())
    return (year_val, month_val, s)

def sort_key(name):
    n = str(name).strip().lower()
    if 'private total' in n: return (1, n)
    if 'life insurance corporation' in n or n == 'lic': return (2, n)
    if 'grand total' in n: return (3, n)
    return (0, n)

def render_styled_df(df, num_cols=[], rupee_cols=[], growth_cols=[], int_cols=[], height=None):
    df_disp = df.copy()
    calculated_height = height if height else max(int((len(df_disp) + 1) * 42) + 25, 160)
    col_config = {}

    for col in rupee_cols:
        if col in df_disp.columns:
            col_config[col] = st.column_config.NumberColumn(label=col, format="₹ %,.2f")
    for col in num_cols:
        if col in df_disp.columns:
            col_config[col] = st.column_config.NumberColumn(label=col, format="%,.2f")
    for col in growth_cols:
        if col in df_disp.columns:
            df_disp[col] = pd.to_numeric(df_disp[col].astype(str).str.replace('%',''), errors='coerce')
            col_config[col] = st.column_config.NumberColumn(label=col, format="%.0f%%")
    for col in int_cols:
        if col in df_disp.columns:
            col_config[col] = st.column_config.NumberColumn(label=col, format="%,d")

    st.dataframe(df_disp, column_config=col_config, use_container_width=True, height=calculated_height, hide_index=True)

# ==============================================================================
# ⚡ CACHED DATA PIPELINE & EXCEL/CSV PARSER
# ==============================================================================
@st.cache_data(show_spinner=False)
def load_and_clean_irdai_cached(file_bytes):
    file_obj = io.BytesIO(file_bytes)
    
    try:
        df_raw = pd.read_excel(file_obj, sheet_name=0, header=None)
    except Exception:
        file_obj.seek(0)
        df_raw = pd.read_csv(file_obj, header=None)

    row_cat = df_raw.iloc[1].astype(str).replace(['nan', 'None', '0', '0.0'], '').str.strip()
    row_sub = df_raw.iloc[2].astype(str).replace(['nan', 'None', '0', '0.0'], '').str.strip()
    row_cat = row_cat.replace('', np.nan).ffill().fillna('')

    raw_cols = []
    for cat, sub in zip(row_cat, row_sub):
        col_name = f"{cat}_{sub}" if (cat and sub and cat != sub) else (cat if cat else sub)
        col_name = (col_name.strip()
                    .replace(' / ', '_').replace('/', '_')
                    .replace(' ', '_').replace(',', '_')
                    .replace('%', 'pct').replace('__', '_'))
        raw_cols.append(col_name)

    seen = {}
    unique_cols = []
    for c in raw_cols:
        if c not in seen:
            seen[c] = 0
            unique_cols.append(c)
        else:
            seen[c] += 1
            unique_cols.append(f"{c}_{seen[c]}")

    df_raw.columns = unique_cols
    df_clean = df_raw.iloc[3:].reset_index(drop=True)

    for col in df_clean.columns:
        if not any(x in str(col).lower() for x in ['sl_no', 'insurer', 'slno', 'sl.no', 'serial']):
            s = df_clean[col].astype(str).str.replace(',', '', regex=False).str.strip()
            df_clean[col] = pd.to_numeric(s, errors='coerce').fillna(0.0)
    return df_clean

@st.cache_data(show_spinner="Processing Analytics Engine...")
def process_data_pipeline(file_bytes):
    df_current = load_and_clean_irdai_cached(file_bytes)
    insurer_cols = [col for col in df_current.columns if 'insurer' in str(col).lower()]
    insurer_col = insurer_cols[0] if insurer_cols else df_current.columns[1]

    try:
        col_p_monthly    = df_current.columns[3]
        col_p_py_ytd     = df_current.columns[5]
        col_p_cy_ytd     = df_current.columns[6]
        col_p_growth_ytd = df_current.columns[7]
        col_p_market_share = df_current.columns[8]
        col_pol_py_ytd   = df_current.columns[12]
        col_pol_cy_ytd   = df_current.columns[13]
        col_lives_py_ytd = df_current.columns[19]
        col_lives_cy_ytd = df_current.columns[20]
        col_sa_monthly   = df_current.columns[24]
        col_sa_cy_ytd    = df_current.columns[27]
    except Exception:
        def find_col(kws, exclude=None):
            for c in df_current.columns:
                s = str(c).lower()
                if all(k in s for k in kws):
                    if exclude and any(e in s for e in exclude): continue
                    return c
            return None
        col_p_cy_ytd = find_col(['up_to'], exclude=['assured','policies','lives']) or df_current.columns[6]
        col_p_py_ytd = find_col(['up_to'], exclude=['assured','policies','lives'])
        col_p_growth_ytd = find_col(['growth']) or df_current.columns[7]
        col_p_market_share = find_col(['market_share']) or df_current.columns[8]
        col_p_monthly = find_col(['for_'], exclude=['assured','policies']) or df_current.columns[3]
        col_pol_cy_ytd = find_col(['up_to','policies'])
        col_pol_py_ytd = find_col(['up_to','policies'])
        col_lives_cy_ytd = find_col(['up_to','lives'])
        col_lives_py_ytd = find_col(['up_to','lives'])
        col_sa_cy_ytd = find_col(['up_to','assured'])
        col_sa_monthly = find_col(['for_','assured'])

    months_list = ['JANUARY','FEBRUARY','MARCH','APRIL','MAY','JUNE','JULY','AUGUST','SEPTEMBER','OCTOBER','NOVEMBER','DECEMBER']
    CURRENT_MONTH_NAME = "Current Month"
    for m in months_list:
        if m.lower() in str(col_p_cy_ytd).lower() or m[:3].lower() in str(col_p_cy_ytd).lower():
            CURRENT_MONTH_NAME = m.title()
            break

    years = [int(m.group()) for c in df_current.columns if (m := re.search(r'20\d{2}', str(c)))]
    current_year = max(years) if years else 2026
    prev_year = current_year - 1

    # CORE TRANSFORMATIONS
    df_work = df_current.copy()
    df_work['segment_type'] = df_work[insurer_col].astype(str).str.strip().map(segment_mapping)
    df_work['insurer_name'] = df_work[insurer_col].where(df_work['segment_type'].isna()).ffill()

    # FYP Matrix
    df_company = df_work[df_work['segment_type'].isna()].copy()
    df_company['Total'] = pd.to_numeric(df_company[col_p_cy_ytd], errors='coerce').fillna(0)
    df_company_totals = (df_company[[insurer_col, 'Total']].rename(columns={insurer_col: 'Insurer'}).dropna(subset=['Insurer']))
    df_segments = df_work.dropna(subset=['segment_type']).copy()
    df_seg_matrix = df_segments[~df_segments['insurer_name'].astype(str).str.lower().str.strip().isin(['grand total','total','industry total'])].copy()
    fyp_seg = df_seg_matrix.pivot_table(index='insurer_name', columns='segment_type', values=col_p_cy_ytd, aggfunc='first', fill_value=0).reset_index()
    for seg in ['ISP','INSP','GSP','GNSP','GYRP']:
        if seg not in fyp_seg.columns: fyp_seg[seg] = 0.0
    fyp_seg = fyp_seg.rename(columns={'insurer_name':'Insurer'})
    fyp_matrix = pd.merge(df_company_totals, fyp_seg, on='Insurer', how='left')
    for seg in ['ISP','INSP','GSP','GNSP','GYRP']: fyp_matrix[seg] = fyp_matrix[seg].fillna(0)
    fyp_matrix = fyp_matrix[['Insurer','Total','ISP','INSP','GSP','GNSP','GYRP']]
    fyp_matrix = fyp_matrix[~fyp_matrix['Insurer'].astype(str).str.lower().str.contains('compiled|refers|statement|changed to|w.e.f|nan|none', na=False)].copy()
    fyp_matrix['sort_col'] = fyp_matrix['Insurer'].apply(sort_key)
    fyp_matrix = fyp_matrix.sort_values('sort_col').drop(columns=['sort_col']).reset_index(drop=True)

    # Segment Summary
    df_seg_summary = df_segments[~df_segments['insurer_name'].astype(str).str.lower().str.strip().isin(['private total','public total','grand total','total','industry total'])].copy()
    summary_matrix = df_seg_summary.pivot_table(index='insurer_name', columns='segment_type', values=col_p_cy_ytd, aggfunc='first', fill_value=0)
    summary_data = [{'Segment': s, 'FYP (in Cr.)': summary_matrix[s].sum() if s in summary_matrix.columns else 0} for s in ['ISP','INSP','GSP','GNSP','GYRP']]
    df_summary = pd.DataFrame(summary_data)
    grand_total = df_summary['FYP (in Cr.)'].sum()
    df_summary['Proportion (%)'] = (df_summary['FYP (in Cr.)'] / grand_total * 100).round(2) if grand_total > 0 else 0.0

    # YoY Growth Table & Monthly Premium Extraction
    matched_rows = []
    matched_monthly_rows = []
    for kw, short_name in company_keywords:
        mask = (df_current[insurer_col].astype(str).str.lower().str.contains(kw, na=False) &
                ~df_current[insurer_col].astype(str).str.lower().str.contains('individual|group|single|non-single|yearly|renewable', na=False))
        sub_df = df_current[mask]
        if not sub_df.empty:
            first_row = sub_df.iloc[0]
            cy_val = float(pd.to_numeric(str(first_row[col_p_cy_ytd]).replace(',', ''), errors='coerce') or 0)
            py_val = float(pd.to_numeric(str(first_row[col_p_py_ytd]).replace(',', ''), errors='coerce') or 0)
            m_val = float(pd.to_numeric(str(first_row[col_p_monthly]).replace(',', ''), errors='coerce') or 0)
            official_growth = pd.to_numeric(str(first_row[col_p_growth_ytd]).replace(',', ''), errors='coerce')
            growth_pct = official_growth if pd.notna(official_growth) else (((cy_val - py_val) / py_val * 100) if py_val > 0 else 0.0)
            
            matched_rows.append({
                'INSURER': short_name,
                f"Up to {CURRENT_MONTH_NAME} {prev_year}": round(py_val, 2),
                f"Up to {CURRENT_MONTH_NAME} {current_year}": round(cy_val, 2),
                'FYP Amount (in Cr)': round(cy_val - py_val, 2),
                'Growth Rate': f"{round(growth_pct):.0f}%",
                'Growth_Num': round(growth_pct)
            })
            matched_monthly_rows.append({
                'INSURER': short_name,
                'Monthly FYP (Cr)': round(m_val, 2)
            })

    final_growth_df = pd.DataFrame(matched_rows).sort_values(by=f"Up to {CURRENT_MONTH_NAME} {current_year}", ascending=False).reset_index(drop=True)
    final_growth_df.insert(0, 'Rank', range(1, len(final_growth_df)+1))
    col_current_ytd = final_growth_df.columns[3]

    df_monthly_company = pd.DataFrame(matched_monthly_rows).sort_values(by='Monthly FYP (Cr)', ascending=False).reset_index(drop=True)

    # Market Share & Split Matrix
    df_segments_ms = df_work.dropna(subset=['segment_type']).copy()
    df_segments_ms['ms_val'] = pd.to_numeric(df_segments_ms[col_p_market_share].astype(str).str.replace(',', '').str.replace('%', ''), errors='coerce').fillna(0)
    ms_piv = df_segments_ms.pivot_table(index='insurer_name', columns='segment_type', values='ms_val', aggfunc='first', fill_value=0)
    for seg in ['ISP','INSP','GSP','GNSP','GYRP']:
        if seg not in ms_piv.columns: ms_piv[seg] = 0.0
    ms_piv['Total'] = ms_piv[['ISP','INSP','GSP','GNSP','GYRP']].sum(axis=1)
    ms_matrix_full = ms_piv.reset_index().rename(columns={'insurer_name':'Insurer'})[['Insurer','Total','ISP','INSP','GSP','GNSP','GYRP']]
    ms_matrix_full = ms_matrix_full[~ms_matrix_full['Insurer'].astype(str).str.lower().str.contains('compiled|refers|statement|nan|none', na=False)]
    ms_matrix_full['sort_col'] = ms_matrix_full['Insurer'].apply(sort_key)
    ms_matrix_full = ms_matrix_full.sort_values('sort_col').drop(columns=['sort_col']).reset_index(drop=True)

    ms_rows = []
    for kw, short_name in company_keywords:
        mask = ms_piv.index.astype(str).str.lower().str.contains(kw, na=False)
        sub = ms_piv[mask]
        if not sub.empty:
            r = sub.iloc[0]
            ind_val, grp_val = r['ISP'] + r['INSP'], r['GSP'] + r['GNSP'] + r['GYRP']
            tot = ind_val + grp_val
            ms_rows.append({
                'INSURER': short_name, 'INDIVIDUAL_VAL': ind_val, 'GROUP_VAL': grp_val, 'Total_VAL': tot,
                'Individual Share (%)': (ind_val/tot*100) if tot>0 else 0.0,
                'Group Share (%)': (grp_val/tot*100) if tot>0 else 0.0
            })
    df_ms = pd.DataFrame(ms_rows)
    df_ms = pd.merge(df_ms, final_growth_df[['INSURER', col_current_ytd, 'FYP Amount (in Cr)']], on='INSURER', how='left')
    df_ms = df_ms.sort_values(by=col_current_ytd, ascending=False).reset_index(drop=True)
    df_ms.insert(0, 'Rank', range(1, len(df_ms)+1))
    group_insurers_list = df_ms[df_ms['Group Share (%)'] > 50]['INSURER'].tolist()

    # Lives Covered & MPH Matrices
    df_sub = df_work.dropna(subset=['segment_type']).copy()
    df_sub['Lives_val'] = pd.to_numeric(df_sub[col_lives_cy_ytd].astype(str).str.replace(',', ''), errors='coerce').fillna(0)
    df_sub['P_val'] = pd.to_numeric(df_sub[col_p_cy_ytd].astype(str).str.replace(',', ''), errors='coerce').fillna(0)
    df_sub['SA_val'] = pd.to_numeric(df_sub[col_sa_cy_ytd].astype(str).str.replace(',', ''), errors='coerce').fillna(0) if col_sa_cy_ytd else 0
    df_sub['Pol_val'] = pd.to_numeric(df_sub[col_pol_cy_ytd].astype(str).str.replace(',', ''), errors='coerce').fillna(0) if col_pol_cy_ytd else 0

    group_lives_matrix = df_sub.pivot_table(index='insurer_name', columns='segment_type', values='Lives_val', aggfunc='first', fill_value=0)
    combined_matrix = pd.DataFrame(index=group_lives_matrix.index)
    for seg in ['ISP','INSP','GSP','GNSP','GYRP']: combined_matrix[seg] = group_lives_matrix.get(seg, 0.0)
    combined_matrix = combined_matrix.reset_index()
    combined_matrix['Total'] = combined_matrix[['ISP','INSP','GSP','GNSP','GYRP']].sum(axis=1)
    combined_matrix = combined_matrix.rename(columns={'insurer_name':'Insurer'})[['Insurer','Total','ISP','INSP','GSP','GNSP','GYRP']]
    mask_junk = (combined_matrix['Insurer'].astype(str).str.lower().str.strip().isin(['grand total','industry total','nan','none','']) |
                 combined_matrix['Insurer'].astype(str).str.lower().str.contains('compiled|refers|statement', na=False))
    df_base = combined_matrix[~mask_junk].copy()
    df_private_total = df_base[df_base['Insurer'].astype(str).str.lower().str.strip()=='private total'].copy()
    df_lic = df_base[df_base['Insurer'].astype(str).str.lower().str.strip().isin(['life insurance corporation of india','lic'])].copy()
    df_individuals = df_base[~(df_base['Insurer'].astype(str).str.lower().str.strip().isin(['private total','lic','life insurance corporation of india']))].copy()
    if not df_private_total.empty: df_private_total['Insurer'] = 'Private Total'
    if not df_lic.empty: df_lic['Insurer'] = 'Life Insurance Corporation of India'
    df_individuals['sort_col'] = df_individuals['Insurer'].apply(sort_key)
    df_individuals = df_individuals.sort_values('sort_col').drop(columns='sort_col').reset_index(drop=True)
    grand_total_values = {
        'Insurer': 'Grand Total',
        'ISP': df_individuals['ISP'].sum() + (df_lic['ISP'].sum() if not df_lic.empty else 0),
        'INSP': df_individuals['INSP'].sum() + (df_lic['INSP'].sum() if not df_lic.empty else 0),
        'GSP': df_individuals['GSP'].sum() + (df_lic['GSP'].sum() if not df_lic.empty else 0),
        'GNSP': df_individuals['GNSP'].sum() + (df_lic['GNSP'].sum() if not df_lic.empty else 0),
        'GYRP': df_individuals['GYRP'].sum() + (df_lic['GYRP'].sum() if not df_lic.empty else 0),
    }
    df_grand_total = pd.DataFrame([grand_total_values])
    df_grand_total['Total'] = df_grand_total[['ISP','INSP','GSP','GNSP','GYRP']].sum(axis=1)
    df_lives_clean = pd.concat([df_individuals, df_private_total, df_lic, df_grand_total], ignore_index=True)

    # Master Policy Schemes Matrix
    mph_piv = df_sub.pivot_table(index='insurer_name', columns='segment_type', values='Pol_val', aggfunc='first', fill_value=0)
    for seg in ['ISP','INSP','GSP','GNSP','GYRP']:
        if seg not in mph_piv.columns: mph_piv[seg] = 0.0
    mph_piv['Total'] = mph_piv[['ISP','INSP','GSP','GNSP','GYRP']].sum(axis=1)
    mph_matrix = mph_piv.reset_index().rename(columns={'insurer_name':'Insurer'})[['Insurer','Total','ISP','INSP','GSP','GNSP','GYRP']]
    mph_matrix = mph_matrix[~mph_matrix['Insurer'].astype(str).str.lower().str.contains('compiled|refers|statement|nan|none', na=False)]
    mph_matrix['sort_col'] = mph_matrix['Insurer'].apply(sort_key)
    mph_matrix = mph_matrix.sort_values('sort_col').drop(columns='sort_col').reset_index(drop=True)

    # Lives Summary
    total_lives_sum = df_grand_total['Total'].values[0] if not df_grand_total.empty else 1.0
    lives_summary_data = []
    for seg in ['ISP','INSP','GSP','GNSP','GYRP']:
        seg_lives = df_grand_total[seg].values[0] if not df_grand_total.empty else 0.0
        prop = (seg_lives / total_lives_sum * 100) if total_lives_sum > 0 else 0.0
        lives_summary_data.append({'Segment': seg, 'Lives Covered': seg_lives, 'Proportion (%)': prop})
    df_lives_summary = pd.DataFrame(lives_summary_data)

    # Growth Combined
    df_seg_growth = df_work.dropna(subset=['segment_type']).copy()
    df_seg_growth['growth_val'] = pd.to_numeric(df_seg_growth[col_p_growth_ytd].astype(str).str.replace(',', ''), errors='coerce')
    seg_matrix = df_seg_growth.pivot_table(index='insurer_name', columns='segment_type', values='growth_val', aggfunc='first').reset_index()
    for seg in ['ISP','INSP','GSP','GNSP','GYRP']:
        if seg not in seg_matrix.columns: seg_matrix[seg] = np.nan
    df_company_rows = df_work[df_work['segment_type'].isna()].copy()
    df_company_rows['Total'] = pd.to_numeric(df_company_rows[col_p_growth_ytd].astype(str).str.replace(',', ''), errors='coerce')
    df_company_totals_g = df_company_rows[[insurer_col,'Total']].rename(columns={insurer_col:'insurer_name'}).dropna(subset=['insurer_name'])
    growth_combined = pd.merge(df_company_totals_g, seg_matrix, on='insurer_name', how='left')
    growth_combined = growth_combined[~growth_combined['insurer_name'].astype(str).str.lower().str.strip().isin(['grand total','total','industry total','nan','none',''])].copy()
    growth_combined = growth_combined[~growth_combined['insurer_name'].astype(str).str.lower().str.contains('compiled|refers|statement', na=False)].copy()
    growth_combined = growth_combined.rename(columns={'insurer_name':'Insurer'})[['Insurer','Total','ISP','INSP','GSP','GNSP','GYRP']]
    growth_combined['sort_col'] = growth_combined['Insurer'].apply(sort_key)
    growth_combined = growth_combined.sort_values('sort_col').drop(columns='sort_col').reset_index(drop=True)

    # Rates Matrix
    df_sub_r = df_work.dropna(subset=['segment_type']).copy()
    df_sub_r['P_val'] = pd.to_numeric(df_sub_r[col_p_monthly].astype(str).str.replace(',', ''), errors='coerce').fillna(0)
    df_sub_r['SA_val'] = pd.to_numeric(df_sub_r[col_sa_monthly].astype(str).str.replace(',', ''), errors='coerce').fillna(0)
    p_matrix = df_sub_r.pivot_table(index='insurer_name', columns='segment_type', values='P_val', aggfunc='first', fill_value=0)
    sa_matrix = df_sub_r.pivot_table(index='insurer_name', columns='segment_type', values='SA_val', aggfunc='first', fill_value=0)
    for seg in ['ISP','INSP','GSP','GNSP','GYRP']:
        if seg not in p_matrix.columns: p_matrix[seg] = 0.0
        if seg not in sa_matrix.columns: sa_matrix[seg] = 0.0
    df_comp = df_work[df_work['segment_type'].isna()].copy()
    df_comp['Total_P'] = pd.to_numeric(df_comp[col_p_monthly].astype(str).str.replace(',', ''), errors='coerce').fillna(0)
    df_comp['Total_SA'] = pd.to_numeric(df_comp[col_sa_monthly].astype(str).str.replace(',', ''), errors='coerce').fillna(0)
    df_comp_totals = df_comp[[insurer_col,'Total_P','Total_SA']].rename(columns={insurer_col:'insurer_name'}).dropna(subset=['insurer_name'])
    rate_list = []
    for _, row in df_comp_totals.iterrows():
        c_name = str(row['insurer_name']).strip()
        tot_rate = (row['Total_P'] / row['Total_SA'] * 1000) if row['Total_SA'] > 0 else np.nan
        rec = {'Insurer': c_name, 'Total': tot_rate}
        for seg in ['ISP','INSP','GSP','GNSP','GYRP']:
            p_val = p_matrix.loc[c_name, seg] if c_name in p_matrix.index else 0.0
            sa_val = sa_matrix.loc[c_name, seg] if c_name in sa_matrix.index else 0.0
            rec[seg] = (p_val / sa_val * 1000) if sa_val > 0 else np.nan
        rate_list.append(rec)
    df_rate_all = pd.DataFrame(rate_list)
    df_rate_all = df_rate_all[~df_rate_all['Insurer'].astype(str).str.lower().str.strip().isin(['nan','none',''])].copy()
    df_rate_all = df_rate_all[~df_rate_all['Insurer'].astype(str).str.lower().str.contains('compiled|refers|statement|changed to|w.e.f', na=False)].copy()
    df_rate_all['sort_col'] = df_rate_all['Insurer'].apply(sort_key)
    df_rate_all = df_rate_all.sort_values('sort_col').drop(columns=['sort_col']).reset_index(drop=True)[['Insurer','Total','ISP','INSP','GSP','GNSP','GYRP']]

    # Average Premium & Sum Assured per Life
    sa_mat = df_sub.pivot_table(index='insurer_name', columns='segment_type', values='SA_val', aggfunc='first', fill_value=0)
    pol_mat = df_sub.pivot_table(index='insurer_name', columns='segment_type', values='Pol_val', aggfunc='first', fill_value=0)
    lives_mat = df_sub.pivot_table(index='insurer_name', columns='segment_type', values='Lives_val', aggfunc='first', fill_value=0)

    avg_sa_rows = []
    for c_name in sa_mat.index.unique():
        rec = {'Insurer': str(c_name).strip()}
        seg_sum = 0.0
        for seg in ['ISP', 'INSP', 'GSP', 'GNSP', 'GYRP']:
            sa_cr = sa_mat.loc[c_name, seg] if seg in sa_mat.columns else 0.0
            count_val = pol_mat.loc[c_name, seg] if seg in ['ISP', 'INSP'] else (lives_mat.loc[c_name, seg] if seg in lives_mat.columns else 0.0)
            avg_val = (sa_cr * 10000000) / count_val if count_val != 0 else np.nan
            if pd.notna(avg_val): seg_sum += avg_val
            rec[seg] = avg_val
        rec['Total'] = seg_sum if seg_sum != 0 else np.nan
        avg_sa_rows.append(rec)
    df_avg_sa = pd.DataFrame(avg_sa_rows)
    df_avg_sa = df_avg_sa[~df_avg_sa['Insurer'].astype(str).str.lower().str.strip().isin(['grand total', 'industry total', 'nan', 'none', ''])].copy()
    df_avg_sa = df_avg_sa[~df_avg_sa['Insurer'].astype(str).str.lower().str.contains('compiled|refers|statement', na=False)].copy()
    df_avg_sa['sort_col'] = df_avg_sa['Insurer'].apply(sort_key)
    df_avg_sa = df_avg_sa.sort_values(by='sort_col').drop(columns='sort_col').reset_index(drop=True)[['Insurer', 'Total', 'ISP', 'INSP', 'GSP', 'GNSP', 'GYRP']]

    p_mat = df_sub.pivot_table(index='insurer_name', columns='segment_type', values='P_val', aggfunc='first', fill_value=0)
    avg_life_rows = []
    for c_name in p_mat.index.unique():
        rec = {'Insurer': str(c_name).strip()}
        seg_sum = 0.0
        for seg in ['ISP', 'INSP', 'GSP', 'GNSP', 'GYRP']:
            prem_cr = p_mat.loc[c_name, seg] if seg in p_mat.columns else 0.0
            count_val = pol_mat.loc[c_name, seg] if seg in ['ISP', 'INSP'] else (lives_mat.loc[c_name, seg] if seg in lives_mat.columns else 0.0)
            avg_val = (prem_cr * 10000000) / count_val if count_val != 0 else np.nan
            if pd.notna(avg_val): seg_sum += avg_val
            rec[seg] = avg_val
        rec['Total'] = seg_sum if seg_sum != 0 else np.nan
        avg_life_rows.append(rec)
    df_avg_life = pd.DataFrame(avg_life_rows)
    df_avg_life = df_avg_life[~df_avg_life['Insurer'].astype(str).str.lower().str.strip().isin(['grand total', 'industry total', 'nan', 'none', ''])].copy()
    df_avg_life = df_avg_life[~df_avg_life['Insurer'].astype(str).str.lower().str.contains('compiled|refers|statement', na=False)].copy()
    df_avg_life['sort_col'] = df_avg_life['Insurer'].apply(sort_key)
    df_avg_life = df_avg_life.sort_values(by='sort_col').drop(columns='sort_col').reset_index(drop=True)[['Insurer', 'Total', 'ISP', 'INSP', 'GSP', 'GNSP', 'GYRP']]

    def get_short_name(full_name):
        fn = str(full_name).lower()
        mapping = {'acko':'Acko','aditya':'Aditya Birla','ageas':'Ageas Federal','aviva':'Aviva','axis':'Axis',
                   'bajaj':'Bajaj','bandhan':'Bandhan','bharti':'Bharti Axa','canara':'Canara HSBC',
                   'creditaccess':'CALI','cali':'CALI','edelweiss':'Edelweiss','generali':'Generali Central',
                   'digit':'Go Digit','hdfc':'HDFC','icici':'ICICI','indiafirst':'IndiaFirst',
                   'indusind':'IndusInd Nippon','kotak':'Kotak','pnb':'PNB MetLife','pramerica':'Pramerica',
                   'sbi':'SBI','shriram':'Shriram','star union':'Star Union','tata':'Tata AIA',
                   'private total':'Private Total','life insurance corporation':'LIC','lic':'LIC'}
        for k,v in mapping.items():
            if k in fn: return v
        return full_name

    df_ann = fyp_matrix[['Insurer','Total']].rename(columns={'Total':'FYP (Cr)'}).copy()
    df_ann = pd.merge(df_ann, df_lives_clean[['Insurer','Total']].rename(columns={'Total':'Lives'}), on='Insurer', how='left')
    df_ann = pd.merge(df_ann, df_rate_all[['Insurer','Total']].rename(columns={'Total':'Premium Rate'}), on='Insurer', how='left')
    df_ann['Insurer'] = df_ann['Insurer'].apply(get_short_name)
    df_ann['sort_col'] = df_ann['Insurer'].apply(sort_key)
    df_ann = df_ann.sort_values('sort_col').drop(columns='sort_col').reset_index(drop=True)

    df_sa_matrix = df_sub.pivot_table(index='insurer_name', columns='segment_type', values='SA_val', aggfunc='first', fill_value=0).reset_index()
    for seg in ['ISP','INSP','GSP','GNSP','GYRP']:
        if seg not in df_sa_matrix.columns: df_sa_matrix[seg] = 0.0
    df_sa_matrix['Total'] = df_sa_matrix[['ISP','INSP','GSP','GNSP','GYRP']].sum(axis=1)
    df_sa_matrix = df_sa_matrix.rename(columns={'insurer_name':'Insurer'})[['Insurer','Total','ISP','INSP','GSP','GNSP','GYRP']]

    return {
        "fyp_matrix": fyp_matrix,
        "df_summary": df_summary,
        "final_growth_df": final_growth_df,
        "df_monthly_company": df_monthly_company,
        "df_ms": df_ms,
        "ms_matrix_full": ms_matrix_full,
        "df_lives_summary": df_lives_summary,
        "df_lives_clean": df_lives_clean,
        "mph_matrix": mph_matrix,
        "growth_combined": growth_combined,
        "df_rate_all": df_rate_all,
        "df_avg_life": df_avg_life,
        "df_avg_sa": df_avg_sa,
        "df_ann": df_ann,
        "df_sa_matrix": df_sa_matrix,
        "CURRENT_MONTH_NAME": CURRENT_MONTH_NAME,
        "current_year": current_year,
        "prev_year": prev_year,
        "col_current_ytd": col_current_ytd,
        "group_insurers_list": group_insurers_list
    }

# ⚡ LAZY EXCEL WORKBOOK GENERATOR
@st.cache_data(show_spinner=False)
def generate_full_excel_lazy(file_bytes):
    data = process_data_pipeline(file_bytes)
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        data["df_summary"].to_excel(writer, sheet_name='Segment Summary', index=False)
        data["fyp_matrix"].to_excel(writer, sheet_name='FYP Matrix', index=False)
        data["final_growth_df"].to_excel(writer, sheet_name='YoY Growth Performance', index=False)
        data["df_ms"].to_excel(writer, sheet_name='Ind vs Group Split', index=False)
        data["ms_matrix_full"].to_excel(writer, sheet_name='Market Share Matrix', index=False)
        data["growth_combined"].to_excel(writer, sheet_name='Premium Growth Matrix', index=False)
        data["df_rate_all"].to_excel(writer, sheet_name='Premium Rates', index=False)
        data["df_lives_summary"].to_excel(writer, sheet_name='Lives Summary', index=False)
        data["df_lives_clean"].to_excel(writer, sheet_name='Lives Matrix', index=False)
        data["mph_matrix"].to_excel(writer, sheet_name='Master Policy Schemes MPH', index=False)
        data["df_avg_life"].to_excel(writer, sheet_name='Average Premium Per Life', index=False)
        data["df_avg_sa"].to_excel(writer, sheet_name='Average Sum Assured', index=False)
        data["df_ann"].to_excel(writer, sheet_name='Executive Annexure', index=False)
    output.seek(0)
    return output.getvalue()

# ==============================================================================
# HEADER BANNER & LOGO
# ==============================================================================
st.markdown("<div style='background: rgba(15, 30, 50, 0.4); border-bottom: 3px solid #e63946; padding: 18px 22px; border-radius: 12px; margin-bottom: 25px;'>", unsafe_allow_html=True)
col_title, col_logo = st.columns([4.8, 1.2])

with col_title:
    st.markdown("""
        <h1 style='color:#ffffff; margin:0 0 8px 0; padding:0; font-size: 44px; font-weight:800; text-shadow: 2px 2px 4px rgba(0,0,0,0.4);'>
            CreditAccess Life New Business Dashboard
        </h1>
        <p style='color:#f1faee; margin:0; padding:0; font-size: 20px;'>
            Suraksha aur Samruddhi | IRDAI Monthly Analytics
        </p>
    """, unsafe_allow_html=True)

with col_logo:
    logo_path = BASE_DIR / "cali_logo.png"
    logo_b64 = get_image_base64(logo_path) if logo_path.exists() else None
    if logo_b64:
        st.markdown(f"""
            <div style='background-color: #ffffff !important; padding: 10px 16px !important; border-radius: 10px !important; border: 2px solid #ffffff !important; text-align: center; box-shadow: 0px 4px 15px rgba(0,0,0,0.4); float: right; max-width: 230px;'>
                <img src="data:image/png;base64,{logo_b64}" style="max-width: 100%; max-height: 75px; width: auto; height: auto; display: block; margin: 0 auto; background-color: #ffffff !important;" />
            </div>
        """, unsafe_allow_html=True)
    else:
        st.markdown("""
            <div style='background-color: #ffffff !important; padding: 10px 16px !important; border-radius: 10px !important; border: 2px solid #ffffff !important; text-align: center; float: right; max-width: 230px;'>
                <h3 style='color:#e63946 !important; margin:0; font-weight:bold;'>CreditAccess Life</h3>
            </div>
        """, unsafe_allow_html=True)
st.markdown("</div>", unsafe_allow_html=True)

# ==============================================================================
# SIDEBAR CONTROLS & BOARDROOM PRESENTATION GENERATOR
# ==============================================================================
st.sidebar.title("Monthly Data Management")

if ENABLE_LOGIN:
    if st.sidebar.button("Logout", use_container_width=True):
        st.session_state["authenticated"] = False
        st.rerun()

SAVE_DIR = Path("saved_months")
SAVE_DIR.mkdir(exist_ok=True)

raw_saved_files = [f.stem for f in SAVE_DIR.glob("*.xlsx")]
saved_files = sorted(raw_saved_files, key=parse_month_year)

selected_saved = st.sidebar.selectbox("Previously Saved Months", ["-- Upload New --"] + saved_files)

uploaded_file_bytes = None
if selected_saved != "-- Upload New --":
    try:
        with open(SAVE_DIR / f"{selected_saved}.xlsx", "rb") as f:
            uploaded_file_bytes = f.read()
        st.sidebar.success(f"Loaded: {selected_saved}")
    except PermissionError:
        st.sidebar.error(f"Permission Denied: '{selected_saved}.xlsx' is currently open in Microsoft Excel. Please close the Excel file and refresh!")
        st.stop()
else:
    raw_upload = st.sidebar.file_uploader("Upload New IRDAI File", type=["xlsx", "csv"])
    if raw_upload is not None:
        uploaded_file_bytes = raw_upload.getbuffer().tobytes() if hasattr(raw_upload, 'getbuffer') else raw_upload.read()

if uploaded_file_bytes is not None:
    current_month_str = datetime.now().strftime("%B_%Y")
    month_name = st.sidebar.text_input("Save as (e.g. January_2026)", value=current_month_str)
    if st.sidebar.button("Save this month"):
        try:
            with open(SAVE_DIR / f"{month_name}.xlsx", "wb") as f:
                f.write(uploaded_file_bytes)
            st.sidebar.success(f"Saved as {month_name}.xlsx")
        except PermissionError:
            st.sidebar.error(f"Permission Denied: '{month_name}.xlsx' is currently open in Microsoft Excel. Please close it first!")

    # RUN CACHED PIPELINE
    data = process_data_pipeline(uploaded_file_bytes)

    fyp_matrix = data["fyp_matrix"]
    df_summary = data["df_summary"]
    final_growth_df = data["final_growth_df"]
    df_monthly_company = data["df_monthly_company"]
    df_ms = data["df_ms"]
    ms_matrix_full = data["ms_matrix_full"]
    df_lives_summary = data["df_lives_summary"]
    df_lives_clean = data["df_lives_clean"]
    mph_matrix = data["mph_matrix"]
    growth_combined = data["growth_combined"]
    df_rate_all = data["df_rate_all"]
    df_avg_life = data["df_avg_life"]
    df_avg_sa = data["df_avg_sa"]
    df_ann = data["df_ann"]
    df_sa_matrix = data["df_sa_matrix"]
    CURRENT_MONTH_NAME = data["CURRENT_MONTH_NAME"]
    current_year = data["current_year"]
    prev_year = data["prev_year"]
    col_current_ytd = data["col_current_ytd"]
    group_insurers_list = data["group_insurers_list"]

    # EXCEL DOWNLOAD BUTTON
    st.sidebar.markdown("<br>", unsafe_allow_html=True)
    st.sidebar.subheader("Export Complete Workbook")
    export_filename = f"{selected_saved if selected_saved != '-- Upload New --' else month_name}_Analytics_Workbook.xlsx"
    
    st.sidebar.download_button(
        label="Download All Tables (.xlsx)",
        data=generate_full_excel_lazy(uploaded_file_bytes),
        file_name=export_filename,
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )

    # BOARDROOM REPORT GENERATOR (PPTX & REAL PDF ONLY)
    st.sidebar.markdown("<br>", unsafe_allow_html=True)
    st.sidebar.subheader("Boardroom Report Generator")
    
    active_m_name = selected_saved if selected_saved != '-- Upload New --' else month_name

    # PPTX DOWNLOAD
    if HAS_PPTX:
        pptx_bytes = generate_boardroom_pptx(data, active_m_name)
        if pptx_bytes:
            st.sidebar.download_button(
                label="Download Boardroom Deck (.pptx)",
                data=pptx_bytes,
                file_name=f"{active_m_name}_Boardroom_Deck.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
    else:
        st.sidebar.info("Run `pip install python-pptx` in terminal to enable PowerPoint (.pptx) downloads.")

    # REAL PDF DOWNLOAD
    if HAS_REPORTLAB:
        pdf_bytes = generate_boardroom_pdf(data, active_m_name)
        if pdf_bytes:
            st.sidebar.download_button(
                label="Download Executive Report (.pdf)",
                data=pdf_bytes,
                file_name=f"{active_m_name}_Executive_Report.pdf",
                mime="application/pdf"
            )
    else:
        st.sidebar.info("Run `pip install reportlab` in terminal to enable PDF (.pdf) downloads.")

    # EXECUTIVE HIGH-LEVEL KPI RIBBON BAR
    tot_ind_fyp = df_summary['FYP (in Cr.)'].sum() if not df_summary.empty else 0.0
    tot_lives_num = df_lives_summary['Lives Covered'].sum() if not df_lives_summary.empty else 0.0
    cali_row_kpi = final_growth_df[final_growth_df['INSURER'].str.upper().isin(['CALI', 'CREDITACCESS', 'CREDITACCESS LIFE'])]
    cali_rank_str = f"#{int(cali_row_kpi['Rank'].values[0])}" if not cali_row_kpi.empty else "N/A"
    cali_fyp_val = f"₹ {cali_row_kpi[col_current_ytd].values[0]:,.2f} Cr" if not cali_row_kpi.empty else "₹ 0.00"

    kpi_col1, kpi_col2, kpi_col3, kpi_col4 = st.columns(4)
    with kpi_col1:
        st.metric("Industry Total FYP", f"₹ {tot_ind_fyp:,.2f} Cr")
    with kpi_col2:
        st.metric("Total Lives Covered", f"{tot_lives_num:,.0f}")
    with kpi_col3:
        st.metric("CALI Market Rank", cali_rank_str)
    with kpi_col4:
        st.metric("CALI YTD FYP", cali_fyp_val)

    st.markdown("<br>", unsafe_allow_html=True)

    # ==============================================================================
    # 👑 TABS
    # ==============================================================================
    tab1, tab2, tab3, tab4, tab5, tab6, tab7, tab8, tab9, tab10 = st.tabs([
        "Market Overview (Fyp Segment-Wise)",
        "Premium Growth & Market Share",
        "CALI Competitor",
        "Group Insurer Comparision",
        "Premium Growth & Rates",
        "Lives Covered & MPHs",
        "Avg SA & Avg Premium",
        "Annexure",
        "Company's Monthly Data",
        "Multi-Month Comparison"
    ])

    # ---------- TAB 1 ----------
    with tab1:
        st.subheader(f"Industry FYP by Segment – {CURRENT_MONTH_NAME}")
        
        col_c1, col_c2, col_c3 = st.columns([1, 2, 1])
        with col_c2:
            df_summary_pie = df_summary[df_summary['FYP (in Cr.)'] >= 1.0]
            fig_fyp = px.pie(
                df_summary_pie, 
                values='FYP (in Cr.)', 
                names='Segment',
                hole=0.35,
                color_discrete_sequence=['#e63946', '#f4a261', '#2a9d8f', "#1e3a5f", '#e9c46a']
            )
            fig_fyp.update_traces(
                textinfo='percent+label', 
                textposition='inside',
                insidetextfont=dict(color='white', size=14, family='Arial'),
                hoverinfo='label+value+percent',
                marker=dict(line=dict(color='#ffffff', width=2))
            )
            
            fig_fyp.update_layout(
                paper_bgcolor='white', plot_bgcolor='white',
                font=dict(color='#1e3a5f', size=14, family='Arial'),
                title=dict(text='FYP Proportion by Segment', x=0.5, y=0.98, xanchor='center', yanchor='top', font=dict(color='#e63946', size=20, family='Arial Black')),
                height=400,
                margin=dict(t=95, b=20, l=20, r=20),
                showlegend=True
            )
            st.plotly_chart(fig_fyp, use_container_width=True)

        render_styled_df(df_summary, rupee_cols=['FYP (in Cr.)'], num_cols=['Proportion (%)'])
        
        st.subheader("Complete FYP Matrix (YTD in Cr.)")
        render_styled_df(fyp_matrix, rupee_cols=['Total','ISP','INSP','GSP','GNSP','GYRP'])
        st.download_button("Download FYP Matrix (.xlsx)", convert_df_to_excel(fyp_matrix, "FYP Matrix"), "fyp_matrix.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

    # ---------- TAB 2 ----------
    with tab2:
        st.subheader("Individual vs Group Portfolio Split")
        
        tier_choice = st.radio("Select View Tier", ["Top 8", "Middle 8", "Bottom 8"], horizontal=True)
        subset = df_ms.iloc[0:8] if tier_choice=="Top 8" else (df_ms.iloc[8:16] if tier_choice=="Middle 8" else df_ms.iloc[16:24])

        names = subset['INSURER'].values[::-1]
        ind = subset['Individual Share (%)'].values[::-1]
        grp = subset['Group Share (%)'].values[::-1]

        fig_split = go.Figure()
        fig_split.add_trace(go.Bar(
            y=names, x=ind, name='Individual Share (%)', orientation='h',
            marker=dict(color='#1e3a5f'),
            text=[f"{v:.0f}%" if v > 5 else "" for v in ind], 
            textposition='inside',
            insidetextanchor='middle',
            textfont=dict(color='white', size=13, family='Arial Black'),
            hovertemplate="<b>%{y}</b><br>Individual Share: %{x:.2f}%<extra></extra>"
        ))
        fig_split.add_trace(go.Bar(
            y=names, x=grp, name='Group Share (%)', orientation='h',
            marker=dict(color='#e63946'),
            text=[f"{v:.0f}%" if v > 5 else "" for v in grp], 
            textposition='inside',
            insidetextanchor='middle',
            textfont=dict(color='white', size=13, family='Arial Black'),
            hovertemplate="<b>%{y}</b><br>Group Share: %{x:.2f}%<extra></extra>"
        ))
        
        fig_split.update_layout(
            barmode='stack',
            paper_bgcolor='white', plot_bgcolor='white',
            font=dict(color='#1e3a5f', size=13),
            title=dict(text="Portfolio Split (% Proportion)", x=0.5, y=0.98, xanchor='center', yanchor='top', font=dict(color='#1e3a5f', size=20, family='Arial Black')),
            legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="center", x=0.5, font=dict(size=13)),
            xaxis=dict(
                title=dict(text="Share (%)", font=dict(color='#1e3a5f')), 
                range=[0, 100], 
                gridcolor='#e0e0e0',
                tickfont=dict(color='#1e3a5f')
            ),
            yaxis=dict(tickfont=dict(color='#002147', size=14, family='Arial')),
            height=460,
            margin=dict(t=95, b=30, l=20, r=20)
        )
        st.plotly_chart(fig_split, use_container_width=True)

        st.subheader("Individual vs Group Percentage Proportions & Values")
        render_styled_df(
            subset[['Rank', 'INSURER', 'INDIVIDUAL_VAL', 'GROUP_VAL', 'Individual Share (%)', 'Group Share (%)']], 
            rupee_cols=['INDIVIDUAL_VAL', 'GROUP_VAL'],
            num_cols=['Individual Share (%)', 'Group Share (%)']
        )

        st.markdown("---")
        st.subheader("Private Insurers YoY YTD Performance")
        display_growth = final_growth_df[['Rank', 'INSURER', f"Up to {CURRENT_MONTH_NAME} {prev_year}",
                                          f"Up to {CURRENT_MONTH_NAME} {current_year}", 'FYP Amount (in Cr)', 'Growth Rate']].copy()
        
        render_styled_df(
            display_growth, 
            rupee_cols=[f"Up to {CURRENT_MONTH_NAME} {prev_year}", f"Up to {CURRENT_MONTH_NAME} {current_year}", 'FYP Amount (in Cr)'],
            growth_cols=['Growth Rate']
        )
        st.download_button("Download Growth Table (.xlsx)", convert_df_to_excel(display_growth, "YoY Growth"), "yoy_growth.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

        st.markdown("---")
        st.subheader("Complete Market Share Matrix (%)")
        render_styled_df(ms_matrix_full, num_cols=['Total','ISP','INSP','GSP','GNSP','GYRP'])

    # ---------- TAB 3 ----------
    with tab3:
        st.subheader("CALI Peer Competitor Analytics")
        growth_ranked = final_growth_df.sort_values(by=col_current_ytd, ascending=False).reset_index(drop=True)
        cali_mask = growth_ranked['INSURER'].str.upper().isin(['CALI','CREDITACCESS','CREDITACCESS LIFE'])
        if cali_mask.any():
            cali_idx = growth_ranked[cali_mask].index[0]
            n = len(growth_ranked)
            start = max(0, n-5) if cali_idx >= n-5 else max(0, cali_idx-2)
            end = n if cali_idx >= n-5 else min(n, cali_idx+3)
            df_peers = growth_ranked.iloc[start:end].copy()
        else:
            df_peers = growth_ranked.tail(5).copy()
            cali_idx = -1
        df_peers = df_peers.sort_values(by=col_current_ytd, ascending=False).reset_index(drop=True)
        st.write(f"**CALI Rank: #{cali_idx+1 if cali_idx>=0 else 'N/A'} of {len(growth_ranked)}**")

        col_peer1, col_peer2 = st.columns(2)

        with col_peer1:
            max_g = max(df_peers['Growth_Num']) if not df_peers.empty else 100
            min_g = min(df_peers['Growth_Num']) if not df_peers.empty else 0
            fig_g = go.Figure()
            fig_g.add_trace(go.Bar(
                x=df_peers['INSURER'], y=df_peers['Growth_Num'],
                marker_color='#ee6c4d',
                text=[f"{v:.1f}%" for v in df_peers['Growth_Num']],
                textposition='outside',
                textfont=dict(color='#0d3b66', size=12, family='Arial Black')
            ))
            fig_g.update_layout(
                title=dict(text="FYP Growth %", font=dict(color='#0d3b66', size=18, family='Arial Black'), x=0.5, y=0.98, xanchor='center', yanchor='top'),
                paper_bgcolor='white', plot_bgcolor='white',
                font=dict(color='#0d3b66', size=12), height=360,
                margin=dict(t=90, b=30, l=20, r=20),
                yaxis=dict(gridcolor='#e0e0e0', ticksuffix='%', range=[min(0, min_g * 1.3), max_g * 1.40])
            )
            st.plotly_chart(fig_g, use_container_width=True)

        with col_peer2:
            max_a = max(df_peers[col_current_ytd]) if not df_peers.empty else 100
            fig_a = go.Figure()
            fig_a.add_trace(go.Bar(
                x=df_peers['INSURER'], y=df_peers[col_current_ytd],
                marker_color='#0d3b66',
                text=[f"₹ {v:,.2f}" for v in df_peers[col_current_ytd]],
                textposition='outside',
                textfont=dict(color='#0d3b66', size=12, family='Arial Black')
            ))
            fig_a.update_layout(
                title=dict(text="FYP Amount (in Crores)", font=dict(color='#0d3b66', size=18, family='Arial Black'), x=0.5, y=0.98, xanchor='center', yanchor='top'),
                paper_bgcolor='white', plot_bgcolor='white',
                font=dict(color='#0d3b66', size=12), height=360,
                margin=dict(t=90, b=30, l=20, r=20),
                yaxis=dict(gridcolor='#e0e0e0', range=[0, max_a * 1.40])
            )
            st.plotly_chart(fig_a, use_container_width=True)

        render_styled_df(
            df_peers[['Rank','INSURER',col_current_ytd,'Growth Rate']].rename(columns={col_current_ytd:'FYP (Cr)'}),
            rupee_cols=['FYP (Cr)'], growth_cols=['Growth Rate']
        )

    # ---------- TAB 4 ----------
    with tab4:
        st.subheader("Group Insurers Tier Analysis (>50% Group)")
        df_grp = final_growth_df[final_growth_df['INSURER'].isin(group_insurers_list)].copy()
        df_grp['FYP_VAL'] = pd.to_numeric(df_grp[col_current_ytd], errors='coerce')
        higher = df_grp[df_grp['FYP_VAL']>=500].sort_values('FYP_VAL', ascending=False)
        middle = df_grp[(df_grp['FYP_VAL']>=100)&(df_grp['FYP_VAL']<500)].sort_values('FYP_VAL', ascending=False)
        bottom = df_grp[df_grp['FYP_VAL']<100].sort_values('FYP_VAL', ascending=False)

        tier_choice = st.radio("Select Tier", ["Higher (≥500 Cr)", "Middle (100-500 Cr)", "Bottom (<100 Cr)"], horizontal=True)
        target = higher if "Higher" in tier_choice else (middle if "Middle" in tier_choice else bottom)

        if not target.empty:
            fig_tier = make_subplots(specs=[[{"secondary_y": True}]])
            
            fig_tier.add_trace(
                go.Bar(
                    x=target['INSURER'], y=target['FYP_VAL'], name='FYP Amount', marker_color='#0d3b66',
                    offsetgroup=1,
                    text=[f"{v:,.1f}" for v in target['FYP_VAL']], textposition='outside',
                    textfont=dict(color='#0d3b66', size=12, family='Arial Black')
                ),
                secondary_y=False
            )
            
            fig_tier.add_trace(
                go.Bar(
                    x=target['INSURER'], y=target['Growth_Num'], name='FYP Growth %', marker_color='#ee6c4d',
                    offsetgroup=2,
                    text=[f"{v:.1f}%" for v in target['Growth_Num']], textposition='outside',
                    textfont=dict(color='#ee6c4d', size=12, family='Arial Black')
                ),
                secondary_y=True
            )

            max_fyp_val = max(target['FYP_VAL']) if not target.empty else 100
            max_growth_val = max(target['Growth_Num']) if not target.empty else 100
            min_growth_val = min(target['Growth_Num']) if not target.empty else 0

            title_text = f"{tier_choice} Market Share"
            
            fig_tier.update_layout(
                barmode='group',
                bargap=0.25,
                bargroupgap=0.1,
                paper_bgcolor='white', plot_bgcolor='white',
                font=dict(color='#0d3b66', size=13), height=460,
                title=dict(text=title_text, font=dict(color='#0d3b66', size=22, family='Arial Black'), x=0.5, y=0.98, xanchor='center', yanchor='top'),
                legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="center", x=0.5, font=dict(size=13)),
                margin=dict(t=95, b=20, l=20, r=20)
            )
            fig_tier.update_yaxes(title_text="FYP Amount (Cr)", secondary_y=False, gridcolor='#e0e0e0', tickfont=dict(color='#0d3b66'), range=[0, max_fyp_val * 1.40])
            fig_tier.update_yaxes(title_text="Growth Rate (%)", secondary_y=True, showgrid=False, tickfont=dict(color='#ee6c4d'), ticksuffix='%', range=[min(0, min_growth_val * 1.3), max_growth_val * 1.40])
            st.plotly_chart(fig_tier, use_container_width=True)

            headers_html = "".join([f"<th style='padding:8px; border:1px solid #d0d0d0; text-align:center;'>{name}</th>" for name in target['INSURER']])
            fyp_html = "".join([f"<td style='padding:8px; border:1px solid #d0d0d0; text-align:center; font-weight:bold;'>{v:,.2f}</td>" for v in target['FYP_VAL']])
            growth_html = "".join([f"<td style='padding:8px; border:1px solid #d0d0d0; text-align:center; font-weight:bold;'>{v:.1f}%</td>" for v in target['Growth_Num']])

            st.markdown(f"""
                <div style='background-color: white; padding: 15px; border-radius: 8px; margin-top: -10px; margin-bottom: 20px; border: 1px solid #e0e0e0;'>
                    <table style='width:100%; border-collapse: collapse; font-family: Arial; font-size: 15px; color: #0d3b66;'>
                        <tr style='background-color: #f8f9fa;'>
                            <th style='padding:8px; border:1px solid #d0d0d0;'></th>
                            {headers_html}
                        </tr>
                        <tr>
                            <td style='padding:8px; border:1px solid #d0d0d0; font-weight:bold; color:#0d3b66;'>FYP Amount (in Cr)</td>
                            {fyp_html}
                        </tr>
                        <tr style='background-color: #fff5f5;'>
                            <td style='padding:8px; border:1px solid #d0d0d0; font-weight:bold; color:#ee6c4d;'>FYP Growth %</td>
                            {growth_html}
                        </tr>
                    </table>
                </div>
            """, unsafe_allow_html=True)

            render_styled_df(
                target[['Rank','INSURER',col_current_ytd,'Growth Rate']].rename(columns={col_current_ytd:'FYP (Cr)'}),
                rupee_cols=['FYP (Cr)'], growth_cols=['Growth Rate']
            )
        else:
            st.warning("No companies in this tier.")

    # ---------- TAB 5 ----------
    with tab5:
        st.subheader("1. YoY Premium Growth Matrix (%)")
        render_styled_df(growth_combined, growth_cols=['Total','ISP','INSP','GSP','GNSP','GYRP'])
        st.download_button("Download Growth Matrix (.xlsx)", convert_df_to_excel(growth_combined, "Growth Matrix"), "growth_matrix.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

        st.markdown("---")
        st.subheader("2. Premium Rates (Premium / SA × 1000)")
        render_styled_df(df_rate_all, rupee_cols=['Total','ISP','INSP','GSP','GNSP','GYRP'])
        st.download_button("Download Premium Rates (.xlsx)", convert_df_to_excel(df_rate_all, "Premium Rates"), "premium_rates.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

    # ---------- TAB 6 ----------
    with tab6:
        st.subheader("1. Segment-Wise Summary of Lives Covered")
        
        col_l1, col_l2, col_l3 = st.columns([1, 2, 1])
        with col_l2:
            df_lives_pie = df_lives_summary[df_lives_summary['Proportion (%)'] >= 0.5]
            fig_lives = px.pie(
                df_lives_pie, 
                values='Lives Covered', 
                names='Segment',
                title=f'Lives Covered Mix ({CURRENT_MONTH_NAME})',
                color_discrete_sequence=['#e63946','#f4a261','#2a9d8f','#1e3a5f','#e9c46a'],
                hole=0.35
            )
            fig_lives.update_traces(
                textinfo='percent+label', 
                textposition='inside',
                insidetextfont=dict(color='white', size=14, family='Arial'),
                hoverinfo='label+value+percent',
                marker=dict(line=dict(color='#ffffff', width=2))
            )
            fig_lives.update_layout(
                paper_bgcolor='white', plot_bgcolor='white',
                font=dict(color='#1e3a5f', size=14, family='Arial'),
                title=dict(text=f'Lives Covered Mix ({CURRENT_MONTH_NAME})', x=0.5, y=0.98, xanchor='center', yanchor='top', font=dict(color='#e63946', size=20, family='Arial Black')),
                height=380,
                margin=dict(t=90, b=20, l=20, r=20),
                showlegend=True
            )
            st.plotly_chart(fig_lives, use_container_width=True)

        st.subheader("Segment-Wise Lives Covered & Proportions")
        render_styled_df(df_lives_summary, int_cols=['Lives Covered'], num_cols=['Proportion (%)'])

        st.markdown("---")
        st.subheader("2. Master Policy Schemes (MPH) Matrix")
        render_styled_df(mph_matrix, int_cols=['Total','ISP','INSP','GSP','GNSP','GYRP'])
        st.download_button("Download MPH Matrix (.xlsx)", convert_df_to_excel(mph_matrix, "MPH Matrix"), "mph_matrix.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

        st.markdown("---")
        st.subheader("3. Complete Lives Covered Matrix")
        render_styled_df(df_lives_clean, int_cols=['Total','ISP','INSP','GSP','GNSP','GYRP'])
        st.download_button("Download Lives Matrix (.xlsx)", convert_df_to_excel(df_lives_clean, "Lives Matrix"), "lives_matrix.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

    # ---------- TAB 7 ----------
    with tab7:
        st.subheader("Average Premium Per Life (₹)")
        render_styled_df(df_avg_life, rupee_cols=['Total','ISP','INSP','GSP','GNSP','GYRP'])
        
        st.subheader("Average Sum Assured Per Life (₹)")
        render_styled_df(df_avg_sa, rupee_cols=['Total','ISP','INSP','GSP','GNSP','GYRP'])

    # ---------- TAB 8 ----------
    with tab8:
        st.subheader("Executive Annexure")
        render_styled_df(df_ann, rupee_cols=['FYP (Cr)', 'Premium Rate'], int_cols=['Lives'])
        st.download_button("Download Executive Annexure (.xlsx)", convert_df_to_excel(df_ann, "Annexure"), "executive_annexure.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

    # ---------- TAB 9 ----------
    with tab9:
        st.subheader("Company Deep Dive Analytics")
        
        selected_insurer_tab9 = st.selectbox("Select Insurer for Deep Dive", ['All Companies'] + sorted(final_growth_df['INSURER'].tolist()), key="tab9_company_select")

        if selected_insurer_tab9 == "All Companies":
            st.info("Select a company from the dropdown above to view comprehensive deep-dive metrics.")
        else:
            target_kw = ""
            for kw, sname in company_keywords:
                if sname == selected_insurer_tab9:
                    target_kw = kw
                    break
            if not target_kw:
                target_kw = selected_insurer_tab9.lower()

            def find_company_row(df, col_name='Insurer'):
                if df.empty or col_name not in df.columns: return pd.Series()
                mask = df[col_name].astype(str).str.lower().str.contains(target_kw, na=False)
                sub = df[mask]
                return sub.iloc[0] if not sub.empty else pd.Series()

            g_row = find_company_row(final_growth_df, 'INSURER')
            ms_row = find_company_row(df_ms, 'INSURER')
            fyp_row = find_company_row(fyp_matrix, 'Insurer')
            rate_row = find_company_row(df_rate_all, 'Insurer')
            growth_row = find_company_row(growth_combined, 'Insurer')
            lives_row = find_company_row(df_lives_clean, 'Insurer')
            mph_row = find_company_row(mph_matrix, 'Insurer')
            sa_row = find_company_row(df_sa_matrix, 'Insurer')
            avg_prem_row = find_company_row(df_avg_life, 'Insurer')
            avg_sa_row = find_company_row(df_avg_sa, 'Insurer')

            c1, c2, c3, c4, c5, c6 = st.columns(6)
            
            rank_val = f"#{int(g_row['Rank'])}" if not g_row.empty and 'Rank' in g_row else "N/A"
            fyp_tot_val = f"₹ {fyp_row['Total']:,.2f}" if not fyp_row.empty and 'Total' in fyp_row else "₹ 0.00"
            growth_val = g_row['Growth Rate'] if not g_row.empty and 'Growth Rate' in g_row else "0%"
            rate_tot_val = f"₹ {rate_row['Total']:,.2f}" if not rate_row.empty and 'Total' in rate_row and pd.notna(rate_row['Total']) else "-"
            sa_tot_val = f"₹ {sa_row['Total']:,.2f}" if not sa_row.empty and 'Total' in sa_row else "₹ 0.00"
            group_share_val = f"{ms_row['Group Share (%)']:.0f}%" if not ms_row.empty and 'Group Share (%)' in ms_row else "0%"

            c1.metric("Rank", rank_val)
            c2.metric("YTD FYP (Cr)", fyp_tot_val)
            c3.metric("YoY Growth", growth_val)
            c4.metric("Premium Rate", rate_tot_val)
            c5.metric("Sum Assured (Cr)", sa_tot_val)
            c6.metric("Group Share", group_share_val)

            st.markdown("---")

            st.subheader(f"Segment Breakdown for {selected_insurer_tab9}")
            segments = ['ISP', 'INSP', 'GSP', 'GNSP', 'GYRP']
            
            seg_table_data = []
            for seg in segments:
                f_val = fyp_row[seg] if not fyp_row.empty and seg in fyp_row else 0.0
                g_val = growth_row[seg] if not growth_row.empty and seg in growth_row else np.nan
                r_val = rate_row[seg] if not rate_row.empty and seg in rate_row else np.nan
                l_val = lives_row[seg] if not lives_row.empty and seg in lives_row else 0
                mp_val = mph_row[seg] if not mph_row.empty and seg in mph_row else 0
                s_val = sa_row[seg] if not sa_row.empty and seg in sa_row else 0.0
                ap_val = avg_prem_row[seg] if not avg_prem_row.empty and seg in avg_prem_row else np.nan
                asa_val = avg_sa_row[seg] if not avg_sa_row.empty and seg in avg_sa_row else np.nan

                seg_table_data.append({
                    'Segment': seg,
                    'FYP (in Cr.)': f_val,
                    'Growth Rate': g_val,
                    'Premium Rate': r_val,
                    'MPHs': mp_val,
                    'Lives Covered': l_val,
                    'Sum Assured (Cr)': s_val,
                    'Avg Premium per Life (₹)': ap_val,
                    'Avg SA per Life (₹)': asa_val
                })

            df_deep_seg = pd.DataFrame(seg_table_data)

            col_t, col_p = st.columns([1.3, 1])
            with col_t:
                render_styled_df(
                    df_deep_seg, 
                    rupee_cols=['FYP (in Cr.)', 'Premium Rate', 'Sum Assured (Cr)', 'Avg Premium per Life (₹)', 'Avg SA per Life (₹)'],
                    growth_cols=['Growth Rate'],
                    int_cols=['MPHs', 'Lives Covered'],
                    height=280
                )

            with col_p:
                df_deep_pie = df_deep_seg[df_deep_seg['FYP (in Cr.)'] >= 0.1]
                if not df_deep_pie.empty:
                    fig_deep = px.pie(
                        df_deep_pie, 
                        values='FYP (in Cr.)', 
                        names='Segment',
                        title=f"{selected_insurer_tab9} Segment Mix",
                        color_discrete_sequence=['#e63946','#f4a261','#2a9d8f','#1e3a5f','#e9c46a'],
                        hole=0.35
                    )
                    fig_deep.update_traces(
                        textinfo='label+percent',
                        textposition='inside',
                        insidetextfont=dict(color='white', size=14, family='Arial'),
                        hoverinfo='label+value+percent',
                        marker=dict(line=dict(color='#ffffff', width=1.5))
                    )
                    fig_deep.update_layout(
                        paper_bgcolor='white', plot_bgcolor='white',
                        font=dict(color='#1e3a5f', size=14, family='Arial'),
                        title=dict(text=f"{selected_insurer_tab9} Segment Mix", x=0.5, y=0.98, xanchor='center', yanchor='top', font=dict(color='#1e3a5f', size=18, family='Arial Black')),
                        height=380, margin=dict(t=90, b=20, l=20, r=20)
                    )
                    st.plotly_chart(fig_deep, use_container_width=True)
                else:
                    st.info("No active FYP segment data available for pie chart visualization.")

    # ---------- TAB 10 (USES FOR THE MONTH MONTHLY PREMIUM DATA) ----------
    with tab10:
        st.subheader("Multi-Month Comparative Analysis")
        all_saved_files = sorted([f for f in SAVE_DIR.glob("*.xlsx")], key=lambda x: parse_month_year(x.stem))
        
        if len(all_saved_files) < 2:
            st.info("To compare performance across periods, please save at least 2 monthly reports in the sidebar!")
        else:
            comp_mode = st.radio(
                "Select Comparison View Mode", 
                ["Industry Tier Comparison (Month 1 vs Month 2)", "Single Company Multi-Month Trend (Jan - Dec)", "Custom Multi-Insurer Peer Comparison"], 
                horizontal=True
            )

            saved_names = [f.stem.replace('_', ' ').title() for f in all_saved_files]

            # --- MODE 1: MONTHLY PREMIUM COMPARISON ---
            if comp_mode == "Industry Tier Comparison (Month 1 vs Month 2)":
                c_m1, c_m2 = st.columns(2)
                with c_m1:
                    month1_str = st.selectbox("Select Base Period (Month 1)", options=saved_names, index=0, key="m1_select_mode1")
                with c_m2:
                    month2_str = st.selectbox("Select Comparison Period (Month 2)", options=saved_names, index=len(saved_names)-1, key="m2_select_mode1")

                if month1_str == month2_str:
                    st.warning("Please select two different months to perform comparative analytics.")
                else:
                    idx1 = saved_names.index(month1_str)
                    idx2 = saved_names.index(month2_str)

                    with open(all_saved_files[idx1], "rb") as f1:
                        data1 = process_data_pipeline(f1.read())
                    with open(all_saved_files[idx2], "rb") as f2:
                        data2 = process_data_pipeline(f2.read())

                    # USES "FOR THE MONTH" MONTHLY PREMIUM DATA
                    g1 = data1["df_monthly_company"][['INSURER', 'Monthly FYP (Cr)']].rename(columns={'Monthly FYP (Cr)': f"Monthly_FYP_{month1_str}"})
                    g2 = data2["df_monthly_company"][['INSURER', 'Monthly FYP (Cr)']].rename(columns={'Monthly FYP (Cr)': f"Monthly_FYP_{month2_str}"})

                    df_comp = pd.merge(g1, g2, on='INSURER', how='outer').fillna(0.0)
                    
                    # DEFINED DIFFERENCE COLUMN FIRST TO PREVENT KEYERROR
                    df_comp['FYP Difference (Cr)'] = df_comp[f"Monthly_FYP_{month2_str}"] - df_comp[f"Monthly_FYP_{month1_str}"]
                    df_comp['Growth Shift (%)'] = np.where(df_comp[f"Monthly_FYP_{month1_str}"] > 0, 
                                                          (df_comp['FYP Difference (Cr)'] / df_comp[f"Monthly_FYP_{month1_str}"]) * 100, 0.0)

                    # HIGHLIGHT CALI SHIFT
                    cali_comp = df_comp[df_comp['INSURER'].str.upper().isin(['CALI', 'CREDITACCESS', 'CREDITACCESS LIFE'])]
                    if not cali_comp.empty:
                        c_row = cali_comp.iloc[0]
                        c_fyp1 = c_row[f"Monthly_FYP_{month1_str}"]
                        c_fyp2 = c_row[f"Monthly_FYP_{month2_str}"]
                        c_delta = c_row['FYP Difference (Cr)']

                        st.markdown(f"""
                            <div style='background: rgba(30, 58, 95, 0.9); border-left: 6px solid #e63946; padding: 16px; border-radius: 10px; margin-bottom: 20px;'>
                                <h3 style='margin:0 0 8px 0; color:#ffffff;'>CALI Monthly Performance Shift: {month1_str} vs {month2_str}</h3>
                                <p style='margin:0; color:#f1faee; font-size: 18px;'>
                                    <b>Monthly FYP Movement:</b> ₹ {c_fyp1:,.2f} Cr ➔ <b>₹ {c_fyp2:,.2f} Cr</b> | 
                                    <b>Difference:</b> <span style='color: {"#2a9d8f" if c_delta>=0 else "#e63946"}; font-weight:bold;'>{'+' if c_delta>=0 else ''}₹ {c_delta:,.2f} Cr</span>
                                </p>
                            </div>
                        """, unsafe_allow_html=True)

                    st.subheader(f"Side-by-Side Monthly FYP Comparison (₹ Cr)")
                    
                    tier_choice_comp = st.radio("Select View Tier for Comparison", ["Top 8", "Middle 8", "Bottom 8"], horizontal=True, key="tier_comp_radio")
                    
                    sorted_comp = df_comp.sort_values(by=f"Monthly_FYP_{month2_str}", ascending=False).reset_index(drop=True)
                    
                    if tier_choice_comp == "Top 8":
                        subset_comp = sorted_comp.iloc[0:8]
                    elif tier_choice_comp == "Middle 8":
                        subset_comp = sorted_comp.iloc[8:16]
                    else:
                        subset_comp = sorted_comp.iloc[16:24]
                    
                    fig_mcomp = go.Figure()
                    fig_mcomp.add_trace(go.Bar(
                        x=subset_comp['INSURER'], y=subset_comp[f"Monthly_FYP_{month1_str}"],
                        name=month1_str, marker_color='#1e3a5f',
                        text=[f"₹{v:,.1f}" for v in subset_comp[f"Monthly_FYP_{month1_str}"]], textposition='outside',
                        textfont=dict(color='#1e3a5f', size=12, family='Arial Black')
                    ))
                    fig_mcomp.add_trace(go.Bar(
                        x=subset_comp['INSURER'], y=subset_comp[f"Monthly_FYP_{month2_str}"],
                        name=month2_str, marker_color='#e63946',
                        text=[f"₹{v:,.1f}" for v in subset_comp[f"Monthly_FYP_{month2_str}"]], textposition='outside',
                        textfont=dict(color='#e63946', size=12, family='Arial Black')
                    ))

                    max_comp_y = max(max(subset_comp[f"Monthly_FYP_{month1_str}"]), max(subset_comp[f"Monthly_FYP_{month2_str}"])) if not subset_comp.empty else 100

                    fig_mcomp.update_layout(
                        barmode='group', paper_bgcolor='white', plot_bgcolor='white',
                        font=dict(color='#1e3a5f', size=13), height=460,
                        title=dict(text=f"{tier_choice_comp} Insurers Monthly FYP: {month1_str} vs {month2_str} (₹ Cr)", x=0.5, y=0.98, xanchor='center', yanchor='top', font=dict(color='#1e3a5f', size=20, family='Arial Black')),
                        legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="center", x=0.5, font=dict(size=13)),
                        margin=dict(t=95, b=30, l=20, r=20),
                        yaxis=dict(gridcolor='#e0e0e0', range=[0, max_comp_y * 1.40])
                    )
                    st.plotly_chart(fig_mcomp, use_container_width=True)

                    st.subheader("Complete Comparative Monthly Difference Matrix")
                    disp_comp = df_comp.sort_values(by=f"Monthly_FYP_{month2_str}", ascending=False).reset_index(drop=True)
                    
                    disp_comp[f"Monthly_FYP_{month1_str}"] = pd.to_numeric(disp_comp[f"Monthly_FYP_{month1_str}"], errors='coerce').fillna(0.0)
                    disp_comp[f"Monthly_FYP_{month2_str}"] = pd.to_numeric(disp_comp[f"Monthly_FYP_{month2_str}"], errors='coerce').fillna(0.0)
                    disp_comp['FYP Difference (Cr)'] = pd.to_numeric(disp_comp['FYP Difference (Cr)'], errors='coerce').fillna(0.0)
                    disp_comp['Growth Shift (%)'] = pd.to_numeric(disp_comp['Growth Shift (%)'], errors='coerce').fillna(0.0)

                    render_styled_df(
                        disp_comp, 
                        rupee_cols=[f"Monthly_FYP_{month1_str}", f"Monthly_FYP_{month2_str}", 'FYP Difference (Cr)'],
                        growth_cols=['Growth Shift (%)']
                    )

            # --- MODE 2: SINGLE COMPANY MONTHLY FYP & YOY GROWTH TREND ---
            elif comp_mode == "Single Company Multi-Month Trend (Jan - Dec)":
                st.subheader("Single Company Historical Performance Trend")
                
                target_company = st.selectbox(
                    "Select Company to Track Across All Saved Months", 
                    sorted(data["final_growth_df"]['INSURER'].tolist()),
                    index=sorted(data["final_growth_df"]['INSURER'].tolist()).index('CALI') if 'CALI' in sorted(data["final_growth_df"]['INSURER'].tolist()) else 0,
                    key="mode2_co_select"
                )

                trend_rows = []
                for filepath in all_saved_files:
                    m_label = filepath.stem.replace('_', ' ').title()
                    with open(filepath, "rb") as f_m:
                        d_m = process_data_pipeline(f_m.read())
                        df_m_comp = d_m["df_monthly_company"]
                        df_g = d_m["final_growth_df"]
                        
                        match_m = df_m_comp[df_m_comp['INSURER'].str.upper() == target_company.upper()]
                        match_g = df_g[df_g['INSURER'].str.upper() == target_company.upper()]
                        
                        m_fyp_val = match_m.iloc[0]['Monthly FYP (Cr)'] if not match_m.empty else 0.0
                        g_rate_val = match_g.iloc[0]['Growth_Num'] if not match_g.empty else 0.0
                        g_rate_str = match_g.iloc[0]['Growth Rate'] if not match_g.empty else "0%"
                        
                        trend_rows.append({
                            'Month': m_label,
                            'Monthly FYP (Cr)': m_fyp_val,
                            'YoY Growth Rate (%)': g_rate_val,
                            'Growth Rate': g_rate_str
                        })

                df_trend = pd.DataFrame(trend_rows)

                if not df_trend.empty:
                    col_tr1, col_tr2 = st.columns(2)
                    
                    with col_tr1:
                        fig_tr_fyp = px.bar(
                            df_trend, x='Month', y='Monthly FYP (Cr)', text='Monthly FYP (Cr)',
                            title=f"{target_company} Monthly FYP Trend Across Months",
                            color_discrete_sequence=['#e63946'] if target_company=='CALI' else ['#1e3a5f']
                        )
                        fig_tr_fyp.update_traces(texttemplate='₹ %{text:,.2f}', textposition='outside')
                        fig_tr_fyp.update_layout(
                            paper_bgcolor='white', plot_bgcolor='white',
                            font=dict(color='#1e3a5f', size=13),
                            title=dict(x=0.5, y=0.92, xanchor='center', yanchor='top', font=dict(color='#1e3a5f', size=18, family='Arial Black')),
                            margin=dict(t=90, b=30, l=20, r=20),
                            yaxis=dict(range=[0, max(df_trend['Monthly FYP (Cr)']) * 1.3 if max(df_trend['Monthly FYP (Cr)'])>0 else 10])
                        )
                        st.plotly_chart(fig_tr_fyp, use_container_width=True)

                    # 🌟 YOY GROWTH RATE TREND CHART
                    with col_tr2:
                        fig_tr_growth = px.line(
                            df_trend, x='Month', y='YoY Growth Rate (%)', text='YoY Growth Rate (%)', markers=True,
                            title=f"{target_company} YoY Growth Rate Trend Over Time",
                            color_discrete_sequence=['#ee6c4d']
                        )
                        fig_tr_growth.update_traces(texttemplate='%{text:.1f}%', textposition='top center', line=dict(width=3), marker=dict(size=10))
                        fig_tr_growth.update_layout(
                            paper_bgcolor='white', plot_bgcolor='white',
                            font=dict(color='#1e3a5f', size=13),
                            title=dict(x=0.5, y=0.92, xanchor='center', yanchor='top', font=dict(color='#1e3a5f', size=18, family='Arial Black')),
                            margin=dict(t=90, b=30, l=20, r=20),
                            yaxis=dict(gridcolor='#e0e0e0', ticksuffix='%')
                        )
                        st.plotly_chart(fig_tr_growth, use_container_width=True)

                    st.subheader(f"Historical Monthly Data Table – {target_company}")
                    render_styled_df(df_trend, rupee_cols=['Monthly FYP (Cr)'], growth_cols=['Growth Rate'])
                    st.download_button("Download Company Trend Data (.xlsx)", convert_df_to_excel(df_trend, f"{target_company}_Trend"), f"{target_company}_Monthly_Trend.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

            # --- MODE 3: CUSTOM MULTI-INSURER PEER COMPARISON ---
            elif comp_mode == "Custom Multi-Insurer Peer Comparison":
                st.subheader("Custom Multi-Insurer Side-by-Side Comparison")

                c_m1, c_m2 = st.columns(2)
                with c_m1:
                    month1_str = st.selectbox("Select Base Period (Month 1)", options=saved_names, index=0, key="m1_select_mode3")
                with c_m2:
                    month2_str = st.selectbox("Select Comparison Period (Month 2)", options=saved_names, index=len(saved_names)-1, key="m2_select_mode3")

                if month1_str == month2_str:
                    st.warning("Please select two different months to perform comparative analytics.")
                else:
                    idx1 = saved_names.index(month1_str)
                    idx2 = saved_names.index(month2_str)

                    with open(all_saved_files[idx1], "rb") as f1:
                        data1 = process_data_pipeline(f1.read())
                    with open(all_saved_files[idx2], "rb") as f2:
                        data2 = process_data_pipeline(f2.read())

                    g1 = data1["df_monthly_company"][['INSURER', 'Monthly FYP (Cr)']].rename(columns={'Monthly FYP (Cr)': f"Monthly_FYP_{month1_str}"})
                    g2 = data2["df_monthly_company"][['INSURER', 'Monthly FYP (Cr)']].rename(columns={'Monthly FYP (Cr)': f"Monthly_FYP_{month2_str}"})

                    df_comp_custom = pd.merge(g1, g2, on='INSURER', how='outer').fillna(0.0)
                    
                    # DEFINED DIFFERENCE COLUMN FIRST TO PREVENT KEYERROR
                    df_comp_custom['FYP Difference (Cr)'] = df_comp_custom[f"Monthly_FYP_{month2_str}"] - df_comp_custom[f"Monthly_FYP_{month1_str}"]
                    df_comp_custom['Growth Shift (%)'] = np.where(df_comp_custom[f"Monthly_FYP_{month1_str}"] > 0, 
                                                          (df_comp_custom['FYP Difference (Cr)'] / df_comp_custom[f"Monthly_FYP_{month1_str}"]) * 100, 0.0)

                    default_peers = [c for c in ['CALI', 'SBI Life', 'HDFC Life', 'ICICI Prudential'] if c in df_comp_custom['INSURER'].tolist()]
                    
                    selected_peers = st.multiselect(
                        "Pick Companies to Compare Side-by-Side:",
                        options=sorted(df_comp_custom['INSURER'].tolist()),
                        default=default_peers,
                        key="custom_peer_multiselect"
                    )

                    if not selected_peers:
                        st.info("Please select at least one insurer from the dropdown above to view comparative analytics.")
                    else:
                        subset_custom = df_comp_custom[df_comp_custom['INSURER'].isin(selected_peers)].sort_values(by=f"Monthly_FYP_{month2_str}", ascending=False).reset_index(drop=True)

                        fig_custom = go.Figure()
                        fig_custom.add_trace(go.Bar(
                            x=subset_custom['INSURER'], y=subset_custom[f"Monthly_FYP_{month1_str}"],
                            name=month1_str, marker_color='#1e3a5f',
                            text=[f"₹{v:,.1f}" for v in subset_custom[f"Monthly_FYP_{month1_str}"]], textposition='outside',
                            textfont=dict(color='#1e3a5f', size=12, family='Arial Black')
                        ))
                        fig_custom.add_trace(go.Bar(
                            x=subset_custom['INSURER'], y=subset_custom[f"Monthly_FYP_{month2_str}"],
                            name=month2_str, marker_color='#e63946',
                            text=[f"₹{v:,.1f}" for v in subset_custom[f"Monthly_FYP_{month2_str}"]], textposition='outside',
                            textfont=dict(color='#e63946', size=12, family='Arial Black')
                        ))

                        max_c_y = max(max(subset_custom[f"Monthly_FYP_{month1_str}"]), max(subset_custom[f"Monthly_FYP_{month2_str}"])) if not subset_custom.empty else 100

                        fig_custom.update_layout(
                            barmode='group', paper_bgcolor='white', plot_bgcolor='white',
                            font=dict(color='#1e3a5f', size=13), height=460,
                            title=dict(text=f"Custom Peers Monthly Comparison: {month1_str} vs {month2_str}", x=0.5, y=0.98, xanchor='center', yanchor='top', font=dict(color='#1e3a5f', size=20, family='Arial Black')),
                            legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="center", x=0.5, font=dict(size=13)),
                            margin=dict(t=95, b=30, l=20, r=20),
                            yaxis=dict(gridcolor='#e0e0e0', range=[0, max_c_y * 1.40])
                        )
                        st.plotly_chart(fig_custom, use_container_width=True)

                        st.subheader("Custom Peer Comparative Difference Matrix")
                        render_styled_df(
                            subset_custom,
                            rupee_cols=[f"Monthly_FYP_{month1_str}", f"Monthly_FYP_{month2_str}", 'FYP Difference (Cr)'],
                            growth_cols=['Growth Shift (%)']
                        )

else:
    st.info("Upload an IRDAI Excel or CSV file or select a previously saved month from the sidebar to begin.")
